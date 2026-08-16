"""DSH ppt-engine runner — stdin JSON → ppt_reflex engine → stdout JSON.

Request shape:
  {"action":"catalog"}
  {"action":"build","template":"academic","style":"academic_rigorous",
   "page_w":960,"page_h":540,"output":"D:/out.pptx",
   "slides":[{"title":"...","archetype":"content","regions":[["r1",60,100,840,360]],
              "elements":[{"id":"t1","type":"title","text":"..."},
                          {"id":"b1","type":"box","text":"...","fill_color":[27,58,92],
                           "shape_id":"rounded_rectangle"},
                          {"id":"img1","type":"image","image_path":"D:/a.png",
                           "layout_mode":"hero_right"}],
              "arrows":[{"from":"t1","to":"b1","text":"flow"}]}]}
"""
import sys, json, io, os, time, traceback, contextlib

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
try:
    sys.stdin.reconfigure(encoding="utf-8")
except Exception:
    pass
sys.path.insert(0, r"D:\ppt")
# Raw stdout buffer captured ONCE — survives redirect_stdout contexts, so
# streaming frames always reach the consumer even while engine prints are
# being diverted to stderr.
_RAW_OUT = sys.stdout.buffer
_FRAMES_OUT = None  # optional frames bridge file (panel preview)

# ── T10 security / persistence constants ──
_BREAKER_STATE = r"D:\ppt\_breaker_state.json"   # CircuitBreaker cross-process persistence
_LIVE_ALLOWED_PREFIXES = ("http://127.0.0.1:8765",)  # live preview whitelist
_IMAGE_MAX_BYTES_RUNNER = 50 * 1024 * 1024           # runner-side image guard (mirrors builder)
# 问题#2: 面板配色 relay — _palette_auto.json（唯一合并点，JS 侧注入将被移除）。
# 路径与 frames_out 同源，可经请求字段 palette_file 覆盖。
_PALETTE_AUTO = r"D:\ppt\_palette_auto.json"
_PALETTE_OVERRIDE_KEYS = ("accent_hex", "bg_hex")

# frames_out atomic-bridge state: we append to a per-process tmp file and
# os.replace() it onto the target at build end — never "truncate + append",
# which interleaves under two concurrent processes.
_CUR_FRAMES_OUT = None
_FRAMES_TMP_PATH = None
_FRAMES_TMP_FILE = None


def _color(c):
    if isinstance(c, list) and len(c) == 3:
        return tuple(int(v) for v in c)
    return None


def _element(b, e):
    t = e.get("type", "text")
    text = str(e.get("text", "") or "")
    region = e.get("region", "main")
    if t == "title":
        spec = b.title(text, region=region)
    elif t == "subtitle":
        spec = b.subtitle(text, region=region)
    elif t == "text":
        spec = b.text(text, style=e.get("style", "Body"), region=region)
    elif t == "bullet":
        spec = b.bullet(text, region=region)
    elif t == "footer":
        spec = b.footer(text)
    elif t == "box":
        spec = b.box(text, style=e.get("style", "Body"), region=region,
                     fill_color=_color(e.get("fill_color")),
                     shape_id=e.get("shape_id", "rounded_rectangle"),
                     ph=e.get("ph"), align_h=e.get("align_h", "left"),
                     recipe=e.get("recipe"))
    elif t == "shape":
        spec = b.shape(e.get("shape_id", "star"), region=region,
                       fill_color=_color(e.get("fill_color")),
                       pw=e.get("pw"), ph=e.get("ph"), text=text)
    elif t == "image":
        spec = b.image(e.get("image_path", ""), region=region, pw=e.get("pw"),
                       ph=e.get("ph"), fit_mode=e.get("fit_mode", "fit"),
                       layout_mode=e.get("layout_mode", ""), caption=e.get("caption", ""))
    elif t == "table":
        headers = e.get("headers") or []
        rows = e.get("rows") or []
        if not isinstance(headers, list):
            headers = []
        if not isinstance(rows, list):
            rows = []
        rows = [r for r in rows if isinstance(r, (list, tuple))]
        spec = b.table(headers, rows, region=region,
                       font_size=e.get("font_size", 12.0),
                       header_bg=_color(e.get("header_bg")))
    elif t == "divider":
        spec = b.divider(region=region, color=_color(e.get("color")))
    else:
        raise ValueError(f"unknown element type: {t}")
    # T10 element id chain: deck id wins over the auto-generated e_N id, so
    # build / frames / inspect all carry the caller's id.
    did = e.get("id")
    if did:
        spec.elem_id = str(did)
    return spec


# ── T9: region-scoped inspection protocol (read-only, no render, no write) ──
# Dialogue layer contract for the agent (quote in the prompt):
#   When the panel writes area/question feedback to _feedback_auto.json, the agent
#   first runs  `python _dsh_ppt_runner.py --inspect <deck> <slide> [elem_id ...]`
#   to cross the fuzzy word with region evidence BEFORE generating options:
#     "挤" (cramped)  → read region.local_density / gap_sequence_std_pt
#     "乱" (messy)    → read region.font_size_levels / alignment_residual_pt / focal_point
#     "丑" (ugly)     → read region.local_color_ratio / pairwise_contrast / hue_harmony
#   Option verbs MUST reuse declare_direction()'s 17 directions (increase_box_height,
#   reduce_text, switch_style, …) — never invent new words. After the user accepts a
#   fix, tag the (inspect output + accepted solution) pair with "region" and feed it
#   to tools/golden_harvest.py.
def _run_inspect(argv: list) -> None:
    """python _dsh_ppt_runner.py --inspect <deck.json> <slide_idx> [elem_id ...]"""
    import contextlib
    try:
        i = argv.index("--inspect")
        deck_path = argv[i + 1]
        slide_idx = int(argv[i + 2])
        elem_ids = argv[i + 3:] or None
    except (ValueError, IndexError):
        print(json.dumps({"ok": False, "runner_error":
                          "usage: --inspect <deck.json> <slide_idx> [elem_id ...]"}, ensure_ascii=False))
        return

    try:
        with open(deck_path, "r", encoding="utf-8") as f:
            req = json.load(f)
        from ppt_reflex.builder import PPTBuilder
        b = PPTBuilder(template=req.get("template", "academic"),
                       style=req.get("style"),
                       overrides=req.get("overrides"),
                       page_w=req.get("page_w", 960), page_h=req.get("page_h", 540),
                       strict_tokens=req.get("strict_tokens", True))
        with contextlib.redirect_stdout(sys.stderr):
            for s in req.get("slides", []):
                by_id = {}
                elements = []
                for e in s.get("elements", []):
                    spec = _element(b, e)
                    elements.append(spec)
                    if e.get("id"):
                        by_id[e["id"]] = spec
                arrows = []
                for a in s.get("arrows", []):
                    frm, to = a.get("from"), a.get("to")
                    if frm in by_id and to in by_id:
                        arrows.append(b.arrow(by_id[frm], by_id[to], text=a.get("text", "")))
                b.add_slide(s.get("title", ""), archetype=s.get("archetype"),
                            params=s.get("params"), regions=s.get("regions"),
                            elements=elements, arrows=arrows,
                            frame=s.get("frame", ""), rail=s.get("rail", ""),
                            corner_mark=s.get("corner_mark", ""))
        result = b.inspect_slide(slide_idx, elem_ids)
        print(json.dumps(result, ensure_ascii=False, default=str))
    except Exception as ex:
        traceback.print_exc()  # debug detail stays on stderr
        print(json.dumps({"ok": False, "runner_error": _structured_error(ex, "inspect_failed")},
                         ensure_ascii=False))


def _safe_render_dir(req):
    """Resolve a render output directory that stays inside the engine workspace."""
    base = os.path.realpath(req.get("cwd") or r"D:\ppt")
    requested = req.get("render_dir") or ""
    if requested and not os.path.isabs(requested):
        requested = os.path.join(base, requested)
    if not requested:
        requested = os.path.join(base, "_render_vision")
    try:
        target = os.path.realpath(requested)
        if os.path.commonpath([base, target]) != base:
            return os.path.join(base, "_render_vision")
        return target
    except Exception:
        return os.path.join(base, "_render_vision")


def _workspace_base(req) -> str:
    """Canonical workspace root for path whitelisting (realpath of cwd)."""
    try:
        return os.path.realpath(req.get("cwd") or os.getcwd())
    except Exception:
        return os.path.realpath(r"D:\ppt")


def _safe_output_path(req) -> tuple[str | None, str | None]:
    """T10: output must be absolute(ized) and inside the workspace; anything else
    returns (None, error) so the caller emits a structured diagnostic and never
    writes outside the sandbox."""
    base = _workspace_base(req)
    out = req.get("output") or ""
    if not out:
        return None, "output path missing"
    if not os.path.isabs(out):
        out = os.path.join(base, out)
    try:
        target = os.path.realpath(out)
        if os.path.commonpath([base, target]) != base:
            return None, f"output path escapes workspace: {out}"
        return target, None
    except Exception as ex:
        return None, f"output path invalid: {ex}"


def _valid_frames_out(req) -> tuple[str | None, str | None]:
    """T10: frames_out must be absolute, inside the workspace and end with .jsonl."""
    base = _workspace_base(req)
    fo = req.get("frames_out") or ""
    if not fo:
        return None, None
    if not os.path.isabs(fo):
        return None, "frames_out must be an absolute path"
    if not fo.lower().endswith(".jsonl"):
        return None, "frames_out must end with .jsonl"
    try:
        target = os.path.realpath(fo)
        if os.path.commonpath([base, target]) != base:
            return None, "frames_out escapes workspace"
        return target, None
    except Exception as ex:
        return None, f"frames_out invalid: {ex}"


def _valid_live_url(req) -> tuple[str | None, str | None]:
    """T10: live preview URL whitelist — only the local panel server is allowed."""
    live = req.get("live") or ""
    if not live:
        return None, None
    if not live.startswith(_LIVE_ALLOWED_PREFIXES):
        return None, ("live URL must start with one of: "
                      + ", ".join(_LIVE_ALLOWED_PREFIXES))
    return live, None


def _merge_palette_overrides(req: dict) -> dict | None:
    """问题#2: 构建前把面板配色 _palette_auto.json 的 accent_hex/bg_hex 合并进
    overrides（agent 与面板对同一 deck 得到同一结果）。幂等：请求显式给出的
    同键优先，palette 不覆盖；文件缺失/损坏时静默跳过（返回 None = 不改 req）。

    这是唯一合并点 —— JS 侧注入将被移除。palette 色属于用户 relay 色，
    合并结果走 overrides 通道，builder 的 strict_tokens 豁免（accent_hex/bg_hex）
    对同一通道生效。
    """
    path = req.get("palette_file") or _PALETTE_AUTO
    try:
        with open(path, "r", encoding="utf-8") as f:
            pal = json.load(f)
        if not isinstance(pal, dict):
            return None
    except Exception:
        return None
    merged = dict(req.get("overrides") or {})
    changed = False
    for key in _PALETTE_OVERRIDE_KEYS:
        if key not in merged and pal.get(key):
            merged[key] = pal[key]
            changed = True
    return merged if changed else None


def _structured_error(ex: Exception, code: str = "build_failed") -> dict:
    """Structured error payload — tracebacks never enter runner_error; debug
    detail goes to stderr instead."""
    msg = str(ex) or type(ex).__name__
    return {"code": code, "message": msg[:300]}


# ── frames_out atomic bridge ────────────────────────────────────────────────

def _open_frames_tmp() -> str | None:
    """Open the per-process tmp file for frames_out. Returns tmp path or None."""
    global _FRAMES_TMP_FILE, _FRAMES_TMP_PATH
    if not _CUR_FRAMES_OUT:
        return None
    _FRAMES_TMP_PATH = f"{_CUR_FRAMES_OUT}.{os.getpid()}.tmp"
    try:
        _FRAMES_TMP_FILE = open(_FRAMES_TMP_PATH, "w", encoding="utf-8")
        return _FRAMES_TMP_PATH
    except OSError:
        _FRAMES_TMP_PATH = None
        return None


def _flush_frames() -> None:
    """Close the tmp file and atomically replace it onto frames_out."""
    global _FRAMES_TMP_FILE, _FRAMES_TMP_PATH
    if _FRAMES_TMP_FILE is not None:
        try:
            _FRAMES_TMP_FILE.close()
        except Exception:
            pass
        _FRAMES_TMP_FILE = None
    if _FRAMES_TMP_PATH and _CUR_FRAMES_OUT:
        try:
            os.replace(_FRAMES_TMP_PATH, _CUR_FRAMES_OUT)
        except OSError:
            pass
    _FRAMES_TMP_PATH = None


# ── CircuitBreaker cross-process persistence ────────────────────────────────

def _load_breaker_state() -> dict:
    try:
        with open(_BREAKER_STATE, "r", encoding="utf-8") as f:
            data = json.load(f)
        if isinstance(data, dict) and isinstance(data.get("decks"), dict):
            return data
    except Exception:
        pass
    return {"decks": {}}


def _save_breaker_state(state: dict) -> None:
    tmp = _BREAKER_STATE + f".{os.getpid()}.tmp"
    try:
        with open(tmp, "w", encoding="utf-8") as f:
            json.dump(state, f, ensure_ascii=False)
        os.replace(tmp, _BREAKER_STATE)
    except OSError:
        try:
            if os.path.exists(tmp):
                os.remove(tmp)
        except OSError:
            pass


def _breaker_before_build(b, req) -> tuple[dict, str]:
    """Load persisted breaker state for this deck fingerprint; wire direction/round.
    Returns (state, deck_fingerprint) for the post-build write-back."""
    from ppt_reflex.design_policy import CircuitBreaker, deck_fingerprint
    fp = deck_fingerprint(b)
    state = _load_breaker_state()
    b._breaker = CircuitBreaker.from_dict(state["decks"].get(fp))
    direction = req.get("direction")
    if direction:
        err = b.declare_direction(direction)
        if err:
            print(f"[runner] direction ignored: {err}", file=sys.stderr)
    state["meta"] = {
        "last_fp": fp,
        "round": req.get("round"),
        "ts": time.time(),
        "cwd": req.get("cwd") or os.getcwd(),
    }
    return state, fp


def _breaker_after_build(state: dict, fp: str, b) -> None:
    """Persist breaker state so build_count accumulates across processes."""
    state["decks"][fp] = b._breaker.to_dict()
    _save_breaker_state(state)



def _render_pngs(b, render_dir):
    """Render each solved slide to PNG for vision-capable models/plugins."""
    import os
    try:
        from PIL import Image, ImageDraw, ImageFont
    except Exception:
        return []
    os.makedirs(render_dir, exist_ok=True)
    from ppt_reflex.grid import GridCanvas, GridConfig, execute_phase1
    paths = []
    for i, spec in enumerate(b._slides):
        try:
            plan = b._plan(spec)
            c = GridCanvas(GridConfig())
            c.checkpoint()
            execute_phase1(plan, c)
            img = Image.new("RGB", (960, 540), "white")
            d = ImageDraw.Draw(img)
            for reg in plan.regions:
                d.rectangle([reg.x, reg.y, reg.x + reg.w, reg.y + reg.h],
                            outline="#94a3b8", width=1)
            for pe in plan.elements:
                p = pe.payload
                x, y, w, h = pe.x, pe.y, pe.w, pe.h
                if p is not None and getattr(p, "fill_color", None):
                    d.rectangle([x, y, x + w, y + h],
                                fill=tuple(int(v) for v in p.fill_color))
                elif p is not None and getattr(p, "shape_id", ""):
                    d.rectangle([x, y, x + w, y + h], outline="#334155", width=1)
                if p is not None and getattr(p, "text", ""):
                    text = str(p.text)
                    font = None
                    try:
                        font = ImageFont.truetype("msyh.ttc",
                                                  max(10, int(getattr(p, "font_size", 14) or 14)))
                    except Exception:
                        font = ImageFont.load_default()
                    color = tuple(int(v) for v in p.font_color) if getattr(p, "font_color", None) else (15, 23, 42)
                    d.multiline_text((x + 6, y + 6), text, fill=color, font=font)
            path = os.path.join(render_dir, "slide_%02d.png" % i)
            img.save(path)
            paths.append({"slide": i, "path": path, "width": 960, "height": 540})
        except Exception:
            continue
    return paths



def main() -> None:
    if "--inspect" in sys.argv:
        _run_inspect(sys.argv)
        return
    try:
        req = json.load(sys.stdin)
    except Exception as ex:
        print(json.dumps({"ok": False, "runner_error": _structured_error(ex, "bad_json")},
                         ensure_ascii=False))
        return

    action = req.get("action", "build")

    # 问题#2: 唯一合并点 — 构建前把面板配色 _palette_auto.json 合并进 overrides
    # （显式 overrides 优先；palette 缺失/损坏静默跳过）。JS 侧注入将被移除，
    # agent 与面板对同一 deck 从同一来源取色，结果一致。
    _palette = _merge_palette_overrides(req)
    if _palette is not None:
        req["overrides"] = _palette

    # ── frames_out: strict whitelist + atomic tmp-file bridge (no truncate+append) ──
    frames_path, frames_err = _valid_frames_out(req)
    if frames_err:
        print(json.dumps({"ok": False, "runner_error": {"code": "frames_out_invalid",
                                                        "message": frames_err}},
                         ensure_ascii=False))
        return
    global _CUR_FRAMES_OUT
    _CUR_FRAMES_OUT = frames_path
    if _CUR_FRAMES_OUT and _open_frames_tmp() is None:
        print(json.dumps({"ok": False, "runner_error": {"code": "frames_out_invalid",
                                                        "message": "cannot open frames tmp file"}},
                         ensure_ascii=False))
        return

    b = None
    breaker_state = None
    breaker_fp = None
    try:
        if action == "catalog":
            from ppt_reflex.grid.templates import list_templates
            from ppt_reflex.builder import load_style_presets, list_archetypes
            presets = load_style_presets()
            styles = []
            for pid, p in presets.get("presets", {}).items():
                c = p.get("color_override", {})
                styles.append({
                    "id": pid, "display_name": p["display_name"], "mood": p["mood"],
                    "theme": p["theme"],
                    "accent": c.get("accent", ""), "bg": c.get("bg", ""),
                })
            print(json.dumps({"ok": True,
                              "templates": list_templates(),
                              "styles": styles,
                              "archetypes": list_archetypes()}, ensure_ascii=False))
            return

        # ── output: must stay inside the workspace — else structured diagnostic ──
        out_path, out_err = _safe_output_path(req)
        if out_err:
            if req.get("stream"):
                _emit_line({"result": {"ok": False, "diagnostics": [
                    {"kind": "output_invalid", "severity": "error", "message": out_err}]}})
            else:
                print(json.dumps({"ok": False, "runner_error": {"code": "output_invalid",
                                                                "message": out_err}},
                                 ensure_ascii=False))
            return

        from ppt_reflex.builder import PPTBuilder
        b = PPTBuilder(template=req.get("template", "academic"),
                       style=req.get("style"),
                       overrides=req.get("overrides"),
                       page_w=req.get("page_w", 960), page_h=req.get("page_h", 540),
                       strict_tokens=req.get("strict_tokens", True))
        breaker_state, breaker_fp = _breaker_before_build(b, req)
        # Engine prints (warnings) go to stderr — stdout carries ONLY the result JSON.
        with contextlib.redirect_stdout(sys.stderr):
            for s in req.get("slides", []):
                regions = s.get("regions")
                by_id = {}
                elements = []
                for e in s.get("elements", []):
                    spec = _element(b, e)
                    elements.append(spec)
                    if e.get("id"):
                        by_id[e["id"]] = spec
                arrows = []
                for a in s.get("arrows", []):
                    frm, to = a.get("from"), a.get("to")
                    # unresolved refs pass through as raw ids — the builder emits
                    # invalid_arrow_ref diagnostics instead of silently dropping
                    arrows.append(b.arrow(by_id.get(frm, frm), by_id.get(to, to),
                                          text=a.get("text", "")))
                b.add_slide(s.get("title", ""), archetype=s.get("archetype"),
                            params=s.get("params"),
                            regions=regions, elements=elements, arrows=arrows,
                            frame=s.get("frame", ""), rail=s.get("rail", ""),
                            corner_mark=s.get("corner_mark", ""))
            live_url, live_err = _valid_live_url(req)
            if live_err:
                print(f"[runner] live push rejected: {live_err}", file=sys.stderr)
            elif live_url:
                _push_preview(b, live_url)
            if req.get("stream"):
                _build_streaming(b, out_path,
                                 _safe_render_dir(req) if req.get("render_png") else None)
                return  # result already emitted as a {"result": ...} JSONL line
            r = b.build(out_path)
        _attach_ascii(b, r)
        if req.get("render_png") or req.get("render_dir"):
            r["rendered_slides"] = _render_pngs(b, _safe_render_dir(req))
        print(json.dumps(r, ensure_ascii=False, default=str))
    except Exception as ex:
        traceback.print_exc()  # debug detail stays on stderr, never in runner_error
        print(json.dumps({"ok": False, "runner_error": _structured_error(ex)},
                         ensure_ascii=False))
    finally:
        _flush_frames()
        if breaker_state is not None and b is not None:
            _breaker_after_build(breaker_state, breaker_fp, b)


def _emit_line(obj) -> None:
    """Write one JSONL record to raw stdout and flush — the streaming channel.
    Also buffers into the frames_out tmp file (atomic os.replace at build end —
    never truncate+append, which interleaves under concurrent processes)."""
    line = json.dumps(obj, ensure_ascii=False, default=str) + "\n"
    _RAW_OUT.write(line.encode("utf-8"))
    _RAW_OUT.flush()
    if _FRAMES_TMP_FILE is not None:
        try:
            _FRAMES_TMP_FILE.write(line)
            _FRAMES_TMP_FILE.flush()
        except Exception:
            pass


def _attach_ascii(b, result: dict) -> None:
    """Append per-slide three-tier ASCII feedback (L0 structure / L1 elements /
    L2 numeric text table) to a build result. Signal elements are marked '?' in L1."""
    from ppt_reflex.grid.ascii_map import render_slide_ascii
    from ppt_reflex.grid import GridCanvas, GridConfig, execute_phase1
    from ppt_reflex.grid.composition import global_composition_check

    pages = []
    try:
        for i, spec in enumerate(b._slides):
            plan = b._plan(spec)
            c = GridCanvas(GridConfig())
            c.checkpoint()
            execute_phase1(plan, c)
            ctx = b._composition_context() if hasattr(b, "_composition_context") else None
            diags = global_composition_check(plan, ctx)
            pages.append(render_slide_ascii(plan, c, i, diagnostics=diags))
    except Exception as ex:
        pages.append({"error": str(ex)})
    result["ascii"] = pages


def _build_streaming(b, output: str, render_dir: str | None = None) -> dict:
    """Build one slide at a time, emitting per-element frames to stdout as JSONL
    ({"frame": {...}} lines), then a final {"result": {...}} line. True
    streaming: the consumer paints elements while the build is still running.

    Failure path still emits a structured {"result": {"ok": false,
    "diagnostics": [...]}} line — never a bare traceback."""
    from ppt_reflex.builder import set_render_frame_hook
    from pptx import Presentation
    from pptx.util import Pt

    def _hex(c):
        return "#%02X%02X%02X" % tuple(int(v) for v in c) if c else None

    def emit(elem_id, ct, payload, x, y, w, h):
        _emit_line({"frame": {
            "slide": _state["slide"], "kind": ct.name.lower(),
            "elem_id": elem_id, "text": (payload.text if payload else "") or "",
            "x": round(x, 1), "y": round(y, 1), "w": round(w, 1), "h": round(h, 1),
            "fill": _hex(payload.fill_color) if payload else None,
            "font_size": (payload.font_size if payload else None) or 0,
        }})

    _state = {"slide": 0}
    set_render_frame_hook(emit)
    try:
        try:
            prs = Presentation()
            prs.slide_width = Pt(b.pw)
            prs.slide_height = Pt(b.ph)
            with contextlib.redirect_stdout(sys.stderr):
                for i in range(len(b._slides)):
                    _state["slide"] = i
                    _emit_line({"frame": {"clear_slide": True, "slide": i}})
                    b.build_single_slide(i, prs=prs)
        finally:
            set_render_frame_hook(None)
        # Second pass off-hook: authoritative build + full diagnostics (frames
        # dedup'd by the consumer because the final result replaces the preview).
        with contextlib.redirect_stdout(sys.stderr):
            r = b.build(output)
        _attach_ascii(b, r)
        if render_dir:
            r["rendered_slides"] = _render_pngs(b, render_dir)
        _emit_line({"result": r})
        return r
    except Exception as ex:
        traceback.print_exc()  # debug detail stays on stderr
        diag = {"kind": "build_crashed", "severity": "error",
                "message": (str(ex) or type(ex).__name__)[:300]}
        _emit_line({"result": {"ok": False, "diagnostics": [diag]}})
        return {"ok": False, "diagnostics": [diag]}


def _push_preview(b, live_url: str) -> None:
    """Solve one slide's layout, push per-element frames to the live server, then
    proceed to the real build. Preview coordinates == final PPTX coordinates
    (same deterministic pipeline)."""
    import urllib.request
    from ppt_reflex.grid import GridCanvas, GridConfig, execute_phase1

    def _hex(c):
        if not c:
            return None
        return "#%02X%02X%02X" % tuple(int(v) for v in c)

    for i, spec in enumerate(b._slides):
        plan = b._plan(spec)
        c = GridCanvas(GridConfig())
        c.checkpoint()
        execute_phase1(plan, c)
        frames = [{"clear_slide": True, "slide": i}]
        for reg in plan.regions:
            frames.append({"slide": i, "kind": "region", "text": reg.region_id,
                           "x": reg.x, "y": reg.y, "w": reg.w, "h": reg.h})
        for pe in plan.elements:
            p = pe.payload
            frames.append({
                "slide": i, "kind": pe.content_type.name.lower(),
                "elem_id": pe.elem_id,
                "text": (p.text if p else "") or "",
                "x": pe.x, "y": pe.y, "w": pe.w, "h": pe.h,
                "fill": _hex(p.fill_color) if p else None,
                "font_size": (p.font_size if p else None) or 0,
            })
        data = json.dumps(frames).encode("utf-8")
        req = urllib.request.Request(live_url.rstrip("/") + "/push",
                                     data=data, headers={"Content-Type": "application/json"})
        try:
            urllib.request.urlopen(req, timeout=3)
        except Exception as ex:
            print(f"[runner] live push failed: {ex}", file=sys.stderr)
            return


if __name__ == "__main__":
    main()

"""
ppt_reflex/builder.py — Sole AI entry point. All engine capabilities exposed here, pure interface zero engine concepts.

from ppt_reflex.builder import PPTBuilder, load_style_presets, save_style_presets, list_style_presets, list_archetypes, get_archetype
from ppt_reflex.grid.templates import list_templates

# Agent workflow: browse lightweight catalogs → user picks template+style → declare
# layout intent only (archetype + params) → element factories carry NO raw color and
# NO hand-written coordinates → build → declare_direction on failure.

print(list_templates())       # [{id, name, description, bg_hex, accent_hex, ...}, ...] — 6 entries
print(list_style_presets())   # [{id, display_name, mood, theme}, ...]                   — 6 entries
print(list_archetypes())      # [{id, name, description, guide}, ...]                    — 12 entries

# Template + style. The engine owns every hue and every coordinate.
builder = PPTBuilder(template="academic", style="academic_rigorous")

builder.add_slide("Cover", archetype="title_cover",
    elements=[
        builder.title("Project Title"),
        builder.subtitle("A one-line tagline"),
    ],
)

builder.add_slide("Grid", archetype="grid_cards", params={"columns": 2, "density": "airy"},
    elements=[builder.box("Feature 1", recipe="card"), builder.box("Feature 2", recipe="card")],
)

result = builder.build("out.pptx")
if not result["ok"]:
    builder.declare_direction("reduce_text")

# ── Escape hatches (advanced, for the human panel only — see end of module) ──
# manual regions + color overrides bypass the token discipline. They break the
# harmony floor and are meant for human panel edits, NOT agent-authored decks.
"""

from __future__ import annotations
import os, tempfile, time, math, json, hashlib
from dataclasses import dataclass, field, replace

from ppt_reflex.grid import (
    GridCanvas, GridConfig, ContentType, ElementPayload, Verdict,
    LayoutPlan, Region, Phase1Element, DecoIntent,
    execute_phase1, execute_phase2, audit_plan, global_composition_check,
    Family, POLICIES, family_of, OverlapVerdict,
    SemanticRole, OVERLAY_ROLES,
)
from ppt_reflex.grid.templates import get_template, TemplateProfile
from ppt_reflex.grid.aesthetics import AestheticsEngine, AestheticViolation, ElemStyle
from ppt_reflex.grid.agent_vocabulary import normalize_fit_mode, reject_unknown_kwargs
from ppt_reflex.grid.archetypes import get_archetype, list_archetypes, get_layout_policy, LayoutPolicy, resolve_archetype
from ppt_reflex.grid.serializer import _render_image, _render_payload  # contain-fit rendering
from ppt_reflex.diff_log import DiffLog  # snapshot-based mutation trace, session lifetime
from ppt_reflex.roundtrip_check import check_overflow  # reopen saved PPTX and verify text fits
from ppt_reflex.color_triangulator import check_slide as tri_check_slide  # bg↔text↔fill color triangle

# ── Paths ──
_PRESETS_PATH = os.path.join(os.path.dirname(__file__), "style_presets.json")


def _ensure_writable_path(path: str) -> str:
    """If the target file is write-locked (open in PowerPoint/WPS), write to a _N
    suffix instead of failing with PermissionError. Returns the actual path used."""
    if not os.path.exists(path):
        return path
    try:
        with open(path, "r+b"):
            return path
    except (PermissionError, OSError):
        stem, ext = os.path.splitext(path)
        n = 1
        while True:
            cand = f"{stem}_{n}{ext}"
            if not os.path.exists(cand):
                return cand
            try:
                with open(cand, "r+b"):
                    n += 1
            except (PermissionError, OSError):
                n += 1


def _find_region_height(plan, pe) -> float | None:
    """Return the height of the region that owns element pe, or None."""
    for region in plan.regions:
        if pe.elem_id in region.elements:
            return region.h
    return None



def load_style_presets() -> dict:
    """Read style presets from disk. Returns full dict with meta + presets. Modify then call save_style_presets()."""
    with open(_PRESETS_PATH, "r", encoding="utf-8") as f:
        return json.load(f)


def save_style_presets(data: dict) -> None:
    """Save style presets to disk. Call load->modify->save."""
    with open(_PRESETS_PATH, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def list_style_presets() -> list[dict]:
    """Lightweight preset list for AI style selection. Each entry: id+display_name+mood+theme. No full color/shape data."""
    data = load_style_presets()
    return [
        {"id": pid, "display_name": p["display_name"], "mood": p["mood"], "theme": p["theme"]}
        for pid, p in data.get("presets", {}).items()
    ]


def _load_single_style_preset(style_id: str) -> dict | None:
    """Load ONLY one style preset from JSON — reads the file but discards all other presets."""
    data = load_style_presets()
    return data.get("presets", {}).get(style_id)


# ── WCAG color (single source: grid/color_utils.py) ──
from ppt_reflex.grid.color_utils import is_dark as _is_dark, hex_to_rgb as _hex_to_rgb, rgb_to_hex as _rgb_to_hex

# ── T8: entry discipline — raw color / raw coordinate tokens are forbidden in agent mode ──
_RAW_COLOR_FORBIDDEN = "raw_color_forbidden"
_RAW_COLOR_OVERRIDE_KEYS = frozenset({
    "bg_hex", "text_hex", "title_hex", "accent_hex", "accent2_hex",
    "gray_hex", "dim_hex", "divider_color_hex",
})
# WCAG / color-triangle contrast violations — a human override that lands here is
# surfaced as human_override_warning (respects the human's decision, never blocks).
_CONTRAST_VIOLATION_KINDS = frozenset({
    "color_contrast", "invisible_text", "text_fill_near_match",
    "dark_bg_dark_text", "light_bg_light_text",
    "tri_bg_text", "tri_bg_fill", "tri_fill_text",
})


def _forbid_raw_color(param: str, fix_hint: str) -> None:
    """Reject an agent-authored raw color token. fix_hint points at the recipe /
    style tier that owns colors. Raises ValueError so add_slide/build surface it."""
    raise ValueError(
        f"{_RAW_COLOR_FORBIDDEN}: '{param}' is a raw color escape hatch — the agent "
        f"must not pick hues. {fix_hint}"
    )


def _has_contrast_violation(diags: list[dict]) -> bool:
    """True when any diagnostic is a WCAG / color-triangle contrast violation."""
    for d in diags:
        kind = d.get("kind", "") or ""
        if kind in _CONTRAST_VIOLATION_KINDS or kind.startswith("tri_"):
            return True
    return False


def _harmony_ok(diags: list[dict]) -> bool:
    """T5: harmony_ok = zero error/warning among the T2–T4 harmony rules.
    Computed on RAW diagnostics (before aggregation) so a trimmed warning can't
    silently flip the result. Signals (advisory) never affect harmony_ok."""
    return not any(
        d.get("harmony") and d.get("severity") in ("error", "warning")
        for d in diags
    )

# ── Style table ──
STYLE = {
    "Heading":    dict(font_size=28, font_bold=True,  font_color=(0x1A,0x1A,0x2E), alignment="CENTER"),
    "Subtitle":   dict(font_size=18, font_color=(0x55,0x55,0x77), alignment="CENTER"),
    "Body":       dict(font_size=14, font_color=(0x33,0x33,0x44), font_name="Microsoft YaHei"),
    "Subheading": dict(font_size=16, font_bold=True,  font_color=(0x1B,0x3A,0x5C)),
    "Caption":    dict(font_size=10, font_color=(0x88,0x88,0x99)),
    "Footer":     dict(font_size=8,  font_color=(0xAA,0xAA,0xBB), alignment="CENTER"),
    "ListItem":   dict(font_size=13, font_color=(0x33,0x33,0x44), font_name="Microsoft YaHei"),
    "Emphasis":   dict(font_size=14, font_bold=True,  font_color=(0xC0,0x39,0x2B)),
}

# ── Built-in shape library ──
SHAPES = {
    "rounded_rectangle": "ROUNDED_RECTANGLE", "rectangle": "RECTANGLE",
    "oval": "OVAL", "parallelogram": "PARALLELOGRAM",
    "diamond": "DIAMOND", "chevron": "CHEVRON",
    "pentagon": "PENTAGON", "hexagon": "HEXAGON",
    "up_arrow": "UP_ARROW", "down_arrow": "DOWN_ARROW",
    "left_arrow": "LEFT_ARROW", "right_arrow": "RIGHT_ARROW",
    "star": "STAR_5_POINT", "triangle": "ISOSCELES_TRIANGLE",
    "home": "HOME_PLATE", "cross": "PLUS",
    "pie": "PIE", "wave": "WAVE", "donut": "DONUT",
    "plaque": "PLAQUE", "sun": "SUN",
}

# ── P0-①: diagnostic aggregation constants ──
_PHASE_ORDER = ["0.5", "1", "2", "2.5", "aesthetics", "3.0", "freeze", "pre", "rt"]
_KIND_BATCH_THRESHOLD = 5   # collapse when >=5 warnings share same kind
_WARN_CAP = 15              # max warnings in aggregated output
_INFO_CAP = 5               # max info diagnostics in aggregated output


# ── Internal spec ──
@dataclass
class _Spec:
    elem_id: str; style: str; text: str = ""; region: str = "main"
    ctype: str = "text"; fill_mode: str = "stack"
    pw: float|None = None; ph: float|None = None
    fill_color: tuple|None = None; shape_id: str = ""
    image_path: str = ""; margin: float = 6.0
    fit_mode: str = "fit"      # fit | fill | crop_center — fit=contain, no crop
    allow_upscale: bool = False # small images stay original size
    layout_mode: str = ""      # hero_top | hero_right | hero_left | center_float | small_inline | grid_2x2 | grid_1x3
    caption: str = ""          # Figure caption text
    # Phase1Element extended params (Fix #6)
    align_h: str = "left"
    allow_shrink: bool = False
    allow_wrap: bool = False
    arrow_slot: float = 48.0
    # shape() 文字参数 override — 不建新 style，直接压过解析值
    font_size_override: float|None = None
    font_color_override: tuple|None = None
    # P1-② table data — passed through _Spec → Phase1Element → render
    table_headers: list[str] = field(default_factory=list)
    table_rows: list[list[str]] = field(default_factory=list)
    # 语义角色声明（"AI 写语义"的接线点）："" = 族扶手默认；
    # entity | connector | annotation | emphasis | backdrop
    # 声明 overlay 类 role 后，碰撞系统放行该元素（如高亮框 emphasis 压在卡片上）
    role: str = ""
    corner_radius: float|None = None    # 显式圆角半径 (pt)，优先于模板/预设默认

@dataclass
class _Arrow:
    deco_id: str; from_elem: str; to_elem: str; text: str = ""
    # T8: color defaults are None = resolve to the template contract color. An
    # explicit color tuple is an agent raw-color token and is rejected in strict mode.
    direction: str = "below"; color: tuple|None = None; width: float = 1.5
    # Fix #8: full DecoIntent params
    margin_pt: float = 8.0
    text_font_size: float = 10.0
    text_color: tuple|None = None
    occlusion_check: bool = True

@dataclass
class _Slide:
    title: str = ""
    regions: list = field(default_factory=list)
    elements: list[_Spec] = field(default_factory=list)
    arrows: list[_Arrow] = field(default_factory=list)
    archetype_id: str = ""  # resolved archetype id for this slide (empty = manual)
    # Slide-level decoration skins — engine solves their coordinates, AI never does.
    frame: str = ""        # "top_bottom_band": thin accent bars framing content top+bottom
    rail: str = ""         # "left" | "right": full-height accent rail
    corner_mark: str = ""  # "tl" | "tr": small corner anchor mark


class PPTBuilder:
    """Sole AI entry point. add_slide -> build. Engine + templates fully transparent.
    No theme layer — template + style + overrides only.
    """

    def __init__(self, template: str = "academic", style: str|None = None,
                 overrides: dict|None = None,
                 page_w: float = 960, page_h: float = 540,
                 template_pptx: str|None = None,
                 strict_tokens: bool = True):
        # T8: strict_tokens=True (default) = agent mode — raw color / coordinate
        # tokens are rejected at the API layer. The human panel passes
        # strict_tokens=False to exercise the overrides/escape hatches.
        self.strict_tokens = strict_tokens
        if strict_tokens and overrides:
            color_keys = sorted(_RAW_COLOR_OVERRIDE_KEYS & set(overrides))
            if color_keys:
                _forbid_raw_color(
                    "overrides={" + ", ".join(f"{k!r}" for k in color_keys) + "}",
                    "Pick a template+style instead; the engine owns every hue. "
                    "bg_hex/accent_hex are advanced escape hatches that break the "
                    "harmony floor — human panel only (strict_tokens=False).",
                )
        # Template: lazy — only this one gets instantiated
        self._t: TemplateProfile = get_template(template)
        self._style_preset: dict|None = None
        self._style_id: str|None = style
        self.pw, self.ph = page_w, page_h
        self._slides: list[_Slide] = []
        self._id = 0
        self._template_pptx = template_pptx
        self._style_body_font: str|None = None
        self._image_layout: dict|None = None
        self._shape_override: dict = {}
        self._caption_n: int = 0          # Figure 编号计数（caption 自动编号）
        self._layout_policy: LayoutPolicy = get_layout_policy(template)
        self.diff_log = DiffLog()
        self._pipeline_cache: dict[int, tuple] = {}

        # P0+P1: circuit breaker — cross-build fix-loop detection
        from ppt_reflex.design_policy import CircuitBreaker
        self._breaker = CircuitBreaker()

        if style:
            self._apply_style(style)

        if overrides:
            self._t = self._t.override(**overrides)

    def set_intent_scope(self, scope: dict):
        """Agent-declared intent scope: {"slide_ids":[2,3], "elem_ids":["box_3"]}.
        Used by DiffLog.scope_alert() to catch say-vs-do mismatch."""
        self.diff_log.set_intent_scope(scope)

    def _apply_style(self, style_id: str) -> None:
        """Load and apply a single style preset — only this one is loaded from JSON into memory."""
        preset = _load_single_style_preset(style_id)
        if not preset:
            return
        self._style_preset = preset
        c = preset["color_override"]
        fo = preset.get("font_override", {})
        overrides = dict(
            bg_hex=c.get("bg", self._t.bg_hex),
            text_hex=c.get("text_primary", self._t.text_hex),
            title_hex=c.get("text_primary", self._t.title_hex),
            accent_hex=c.get("accent", self._t.accent_hex),
            accent2_hex=c.get("warn", self._t.accent2_hex),
            gray_hex=c.get("text_secondary", self._t.gray_hex),
            dim_hex=c.get("surface", self._t.dim_hex),
            title_size=fo.get("scale_h1", self._t.title_size),
            body_size=fo.get("scale_body", self._t.body_size),
            subtitle_size=fo.get("scale_h2", self._t.subtitle_size),
            divider_color_hex=c.get("accent", self._t.divider_color_hex),
        )
        self._t = self._t.override(**overrides)
        self._style_body_font = fo.get("body_font")
        self._image_layout = preset.get("image_layout", None)
        self._shape_override = preset.get("shape_override", {}) or {}

    # ── slide ──
    def add_slide(self, title: str = "", *, archetype: str|None = None,
                  params: dict|None = None,
                  regions: list|None = None,
                  elements: list|None = None, arrows: list|None = None,
                  frame: str = "", rail: str = "", corner_mark: str = "") -> int:
        """Add a slide. If archetype is given, auto-resolve regions + auto-route elements via zone_map.

        Manual regions override archetype. Auto-routing only applies to elements without explicit region.
        params: archetype layout parameters (e.g. {"columns": 3, "density": "airy"}) — supported by
        grid_cards; the engine solves coordinates, the caller never writes them.
        frame/rail/corner_mark: slide-level decoration skins (engine-solved geometry).
        """
        arch = None
        resolved_regions = regions

        if archetype:
            try:
                arch = resolve_archetype(archetype, params or None)
            except (KeyError, ValueError):
                # Unknown archetype or unsupported params — fall back to base archetype
                try:
                    arch = get_archetype(archetype)
                except KeyError:
                    pass

        if arch and resolved_regions is None:
            # LayoutPolicy.content_inset 作为 Region 内边距传入（第 7 元），
            # 不再预先收缩坐标——旧版"预收缩 + Region 默认 inset"双重内缩（2026-08 审查）
            inset = getattr(self._layout_policy, 'content_inset', 12)
            resolved_regions = []
            for r in arch.regions:
                name, x, y, w, h, z = r
                resolved_regions.append((name, x, y, w, h, z, inset))

        if resolved_regions is None:
            resolved_regions = [("main", 60, 60, 840, 420, 1)]

        # Auto-route elements via archetype zone_map + distribute（分布组轮流分配）
        routed_elements = []
        if arch and arch.zone_map and elements:
            _dist_counters: dict[str, int] = {}
            zone_names = [r[0] for r in resolved_regions]
            for e in elements:
                if e.region and e.region != "main":
                    routed_elements.append(e)  # explicit region — skip auto-route
                    continue
                # Determine element type from _Spec fields
                etype = self._elem_type(e)
                # 分布组优先：同类元素按声明顺序轮流分配（grid_cards 4 卡位 / comparison A-B）
                dist = getattr(arch, "distribute", {}).get(etype)
                if dist:
                    available = [z for z in dist if z in zone_names]
                    if available:
                        n = _dist_counters.get(etype, 0)
                        routed_elements.append(replace(e, region=available[n % len(available)]))
                        _dist_counters[etype] = n + 1
                        continue
                target_zone = arch.zone_map.get(etype)
                if target_zone:
                    # Check if zone exists in resolved regions
                    if target_zone in zone_names:
                        # Copy, don't mutate caller's _Spec
                        routed_elements.append(replace(e, region=target_zone))
                        continue
                routed_elements.append(e)
        elif elements:
            routed_elements = list(elements)

        self._slides.append(_Slide(title, resolved_regions, routed_elements, arrows or [],
                                   archetype_id=arch.id if arch else "",
                                   frame=frame, rail=rail, corner_mark=corner_mark))
        self._pipeline_cache.pop(len(self._slides) - 1, None)  # invalidate new slide slot
        self._warn_small_header(resolved_regions, routed_elements)
        return len(self._slides) - 1

    def _warn_small_header(self, regions: list, elements: list) -> None:
        """P1: if a title lives in a region too short for its font, warn at declaration
        time instead of failing at freeze. Saves a build cycle."""
        try:
            title_specs = [e for e in elements
                           if getattr(e, 'style', '') in ('Heading', 'Subheading')]
            if not title_specs or not regions:
                return
            for e in title_specs:
                region = next((r for r in regions if r[0] == e.region), None)
                if not region:
                    continue
                region_h = region[4]
                fs = getattr(e, 'font_size', None) or self._title_pt_for(e.style)
                ls = 1.4
                min_h = fs * ls * 1.15 + 10
                if region_h < min_h:
                    print(f"[PPTBuilder] WARN: region '{e.region}' height {region_h}pt < "
                          f"title needs ~{min_h:.0f}pt (font {fs}pt × {ls} line-spacing). "
                          f"Bump the region height before build to avoid overflow.")
        except Exception:
            pass

    def _title_pt_for(self, style: str) -> float:
        return 28.0

    @staticmethod
    def _elem_type(e: _Spec) -> str:
        """Map _Spec → zone_map key. Uses ctype + style for accurate routing.

        Priority: ctype (deterministic) > style (hint). b.title()/bullet()/etc.
        all return ctype="text" — their style distinguishes them."""
        # ctype-first: these are unambiguous
        if e.ctype == "textbox":   return "box"
        if e.ctype == "shape":     return "shape"
        if e.ctype == "image":     return "image"
        if e.ctype == "table":     return "table"
        if e.ctype == "footer":    return "footer"

        # ctype="text" + style disambiguation
        if e.style == "Heading":   return "title"
        if e.style == "Subtitle":  return "subtitle"
        if e.style == "ListItem":  return "bullet"
        if e.style == "Caption":   return "text"
        # "Body"/"Emphasis"/"Subheading" → treated as generic text
        return "text"

    def fix_slide(self, slide_idx: int, title: str|None = None, *,
                  archetype: str|None = None,
                  regions: list|None = None,
                  elements: list|None = None,
                  arrows: list|None = None) -> int:
        """P1-①: modify an existing slide in-place. Returns slide_idx on success, -1 on bad index.

        Only the passed keyword arguments are changed — omitted args keep current values.
        If archetype is passed, auto-resolves regions + auto-routes elements (like add_slide)."""
        if slide_idx < 0 or slide_idx >= len(self._slides):
            return -1
        s = self._slides[slide_idx]
        if title is not None:
            s.title = title
        if archetype is not None:
            arch = get_archetype(archetype)
            s.archetype_id = arch.id
            s.regions = [(r[0], r[1], r[2], r[3], r[4],
                          r[5] if len(r) > 5 else slide_idx + 1)
                         for r in arch.regions]
            if elements is not None:
                routed = []
                for e in elements:
                    if e.region and e.region != "main":
                        routed.append(e)
                        continue
                    etype = PPTBuilder._elem_type(e)
                    target = arch.zone_map.get(etype)
                    if target and target in [r[0] for r in s.regions]:
                        routed.append(replace(e, region=target))
                    else:
                        routed.append(e)
                s.elements = routed
        if elements is not None and archetype is None:
            s.elements = elements
        if regions is not None:
            s.regions = [(r[0], r[1], r[2], r[3], r[4],
                          r[5] if len(r) > 5 else slide_idx + 1)
                         for r in regions]
        if arrows is not None:
            s.arrows = arrows
        self._pipeline_cache.pop(slide_idx, None)
        return slide_idx

    def _spec_hash(self, spec: _Slide) -> str:
        """P1-①: structural fingerprint — same content → same hash, regardless of elem_id."""
        data = {"title": spec.title, "regions": spec.regions, "n_arrows": len(spec.arrows),
                "frame": spec.frame, "rail": spec.rail, "corner_mark": spec.corner_mark}
        els = []
        for e in spec.elements:
            d = {"ctype": e.ctype, "text": e.text, "style": e.style, "region": e.region,
                 "pw": e.pw, "ph": e.ph, "fill_color": e.fill_color, "shape_id": e.shape_id,
                 "image_path": e.image_path, "layout_mode": e.layout_mode,
                 "table_headers": e.table_headers, "table_rows_len": len(e.table_rows) if e.table_rows else 0,
                 "align_h": e.align_h, "fill_mode": e.fill_mode,
                 "corner_radius": e.corner_radius}
            els.append(d)
        data["elements"] = els
        return hashlib.md5(json.dumps(data, sort_keys=True, default=str).encode()).hexdigest()

    def _pipeline_stages(self, i: int, plan, c, diags: list) -> None:
        """Per-slide pipeline stages after plan.validate + execute_phase1:
        audit → collision → phase2 → composition → aesthetics → color triangle → freeze → pre_commit.

        2026-08 审查：这套阶段曾复制 3 份（build_single_slide/build/rebuild）且已漂移
        （build_single_slide 漏掉颜色三角，rebuild 漏 freeze 几何上下文）。现在单源。
        """
        audit_plan(plan, c)
        for d in plan.diagnostics:
            diags.append(_diag(i, "1", d))

        # Phase 1.5: activate dormant three-layer collision system
        diags.extend(_activate_collision(plan, c, i))

        # Phase 2: decoration layer resolution
        decos = execute_phase2(plan, c)
        for d in decos:
            if d.deco_type == "arrow" and d.x2:
                c.register_decoration(d.deco_id, "arrow", d.x1, d.y1, d.x2, d.y2,
                    line_color=d.style.get("line_color", (0x66,0x66,0x66)),
                    line_width_pt=d.style.get("line_width_pt", 1.5),
                    text=d.text, font_size=d.text_font_size, font_color=d.text_color)
            for w in d.occlusion_warnings:
                diags.append(_diag(i, "2", None, kind="arrow_occlusion", severity="warning",
                                   deco_id=d.deco_id, message=w))

        # Phase 2.5: global composition check (spatial + typography + T2–T4 harmony)
        for ci in global_composition_check(plan, self._composition_context()):
            ci.setdefault("kind", ci.get("category", "composition"))
            diags.append(_diag(i, "2.5", ci))

        # Aesthetics (WCAG floor)
        diags.extend(self._run_aesthetics(c, plan))

        # Color triangle: bg ↔ text ↔ fill constraint system (Phase 3.0)
        tri_elems = []
        for pe in plan.elements:
            p = pe.payload
            if not p or not p.text.strip():
                continue
            tri_elems.append({
                "elem_id": pe.elem_id,
                "font_size": p.font_size,
                "font_bold": p.font_bold,
                "font_color_rgb": p.font_color,
                "fill_color_rgb": p.fill_color,
            })
        tri_issues = tri_check_slide(tri_elems, self._t)
        for ti in tri_issues:
            diags.append({
                "slide": i, "phase": "3.0", "kind": f"tri_{ti.edge.replace('↔','_')}",
                "severity": ti.level, "elem_id": ti.elem_id,
                "message": ti.message,
            })

        # Freeze: 2D overflow check on locked coordinates (same data renderer will use)
        from ppt_reflex.grid.text_metrics import check_overflow_2d
        for pe in plan.elements:
            p = pe.payload
            if not p or not p.text.strip():
                continue
            if pe.content_type not in (ContentType.TEXT, ContentType.TEXTBOX):
                continue
            v_auto_fit = not pe.height_is_locked
            h_auto_fit = not pe.width_is_locked
            issues = check_overflow_2d(
                p.text, p.font_size,
                box_w=pe.w, box_h=pe.h,
                line_spacing=p.line_spacing,
                v_auto_fit=v_auto_fit,
                h_auto_fit=h_auto_fit,
            )
            for iss in issues:
                severity = "warning" if pe.content_type == ContentType.TEXTBOX else iss["level"]
                region_h = _find_region_height(plan, pe)
                diags.append({
                    "slide": i, "phase": "freeze",
                    "kind": iss["kind"],
                    "severity": severity,
                    "elem_id": pe.elem_id,
                    "message": iss["message"],
                    "text_height_pt": round(iss.get("rendered_h", 0), 1),
                    "box_height_pt": round(iss.get("box_h", 0), 1),
                    "overflow_pt": round(iss.get("overflow_pt", 0), 1),
                    "region_height_pt": region_h,
                    "recommended_region_h": round((iss.get("rendered_h", 0) + 12), 1) if iss.get("rendered_h") else None,
                    "font_size": iss.get("font_size"),
                    "line_count": iss.get("line_count"),
                    "options": iss.get("options", []),
                })

        # pre_commit_validation — bounds/overflow/role conflicts
        pv = c.pre_commit_validation()
        for err in pv.get("errors", []):
            diags.append(_diag(i, "pre", None, kind="validation_error", severity="error",
                               elem_id=err.get("owner_id",""), message=err.get("detail","")))
        for warn in pv.get("warnings", []):
            diags.append(_diag(i, "pre", None, kind="validation_warning", severity="warning",
                               elem_id=warn.get("owner_id",""), message=warn.get("detail","")))
        for adv in pv.get("advisories", []):
            diags.append(_diag(i, "pre", None, kind="advisory", severity="info",
                               elem_id=adv.get("owner_id",""), message=adv.get("detail","")))

    def build_single_slide(self, slide_idx: int, prs=None) -> dict:
        """Build exactly one slide (for DeckPlanner harness). Returns per-slide diagnostics."""
        if slide_idx < 0 or slide_idx >= len(self._slides):
            return {"ok": False, "diagnostics": [{"kind": "invalid_index", "severity": "error",
                     "message": f"slide_idx {slide_idx} out of range [0, {len(self._slides)})"}]}
        spec = self._slides[slide_idx]
        plan = self._plan(spec)
        c = GridCanvas(GridConfig())
        c.checkpoint()

        self.diff_log.roll()
        diags: list[dict] = []

        plan.validate(verbose=False)
        for d in plan.diagnostics:
            diags.append(_diag(slide_idx, "0.5", d))

        execute_phase1(plan, c)
        self._pipeline_stages(slide_idx, plan, c, diags)

        if prs is not None:
            cap_state = {"n": 0, "format": self.caption_format()}
            _render_slide(prs, c, self._t, slide_index=slide_idx, total_slides=len(self._slides),
                          caption_state=cap_state, slide_spec=spec)
        self.diff_log.snap_after(plan, slide_idx)

        errs = [d for d in diags if d.get("severity") in ("error",)]
        warns = [d for d in diags if d.get("severity") in ("warning","warn")]
        return {"ok": len(errs) == 0, "geometry_ok": len(errs) == 0,
                "harmony_ok": _harmony_ok(diags),
                "diagnostics": diags,
                "summary": f"slide {slide_idx}: {len(diags)} issues ({len(errs)} errors, {len(warns)} warnings)",
                "roundtrip_ok": True}  # single-slide mode defers roundtrip to deck-level build()

    def inspect_slide(self, slide_idx: int, elem_ids: list | None = None) -> dict:
        """Activated three-layer view of one slide — Supply (agent view), Spatial
        (geometry), Profile (layout inference). Memory-only, no render.

        Returns a dict an AI can act on without opening the PPTX:
          - supply:  L0 slide overview (zones, free rects, density)
          - spatial: nearest_neighbor / alignment_groups / gap matrix
          - profile: inferred decorative/title/footer roles
          - overlap: cross-region overlap diagnostics
          - violations / signals: T5 dual-channel harmony diagnostics (scope="slide")

        T9: pass elem_ids to scope the output to those elements + their local
        neighborhood (nearest neighbors + same-region elements). In scoped mode the
        spatial/supply views are cropped and a `region` block adds local metrics:
        local density vs page mean, in-region alignment residual, gap-sequence rhythm,
        in-region font-size hierarchy levels, pairwise contrast, local color ratio.
        Memory-only throughout — never renders, never writes a PPTX.
        """
        if slide_idx < 0 or slide_idx >= len(self._slides):
            return {"ok": False, "error": f"slide_idx {slide_idx} out of range [0, {len(self._slides)})"}
        spec = self._slides[slide_idx]
        plan = self._plan(spec)
        c = GridCanvas(GridConfig())
        c.checkpoint()
        execute_phase1(plan, c)
        audit_plan(plan, c)
        overlap_diags = _activate_collision(plan, c, slide_idx)

        from ppt_reflex.grid.supply import Supply
        from ppt_reflex.grid.spatial import SpatialIndex
        from ppt_reflex.grid.profiles import infer_profile
        from ppt_reflex.grid.composition import global_composition_check

        supply = Supply()
        l0 = supply.level0(c.info_grid)
        spatial = SpatialIndex()
        spatial.rebuild(c.info_grid)
        profile = infer_profile(c.info_grid)
        diags = global_composition_check(plan, self._composition_context())
        violations = [d for d in diags if d.get("channel") == "violation"]
        signals = [d for d in diags if d.get("channel") == "signal"]

        base = {
            "ok": True, "slide": slide_idx,
            "profile": {
                "name": profile.name,
                "decorative_elements": sorted(profile.decorative_elements),
            },
            "overlap": overlap_diags,
        }

        if not elem_ids:
            base.update({
                "scope": "slide",
                "supply": l0,
                "spatial": {
                    "nearest_neighbor": spatial.nearest_neighbor,
                    "alignment_groups": spatial.alignment_groups,
                    "gap_matrix_rows": spatial.gap_matrix_rows,
                    "gap_matrix_cols": spatial.gap_matrix_cols,
                    "density_heatmap": spatial.density_heatmap,
                    "orphans": sorted(spatial.orphans),
                },
                "violations": violations,
                "signals": signals,
            })
            return base

        scoped = self._scope_expand(plan, spatial, set(elem_ids))
        nn = {k: v for k, v in spatial.nearest_neighbor.items() if k in scoped}
        orphans = sorted(set(spatial.orphans) & scoped)
        base.update({
            "scope": "region",
            "elem_ids": sorted(scoped),
            "supply": {"scope": "region", "element_ids": sorted(scoped)},
            "spatial": {"nearest_neighbor": nn, "orphans": orphans,
                        "alignment_groups": {str(k): [i for i in v if i in scoped]
                                             for k, v in spatial.alignment_groups.items()}},
            "region": self._region_metrics(plan, scoped),
            "violations": [v for v in violations if not v.get("elem_id") or v["elem_id"] in scoped],
            "signals": [s for s in signals if not s.get("elem_id") or s["elem_id"] in scoped],
        })
        return base

    def _scope_expand(self, plan, spatial, eids: set) -> set:
        """Expand a seed elem set to its local neighborhood: nearest neighbors +
        every element sharing a region with a seed."""
        scoped = set(eids)
        for eid in list(eids):
            nn = spatial.nearest_neighbor.get(eid)
            if nn:
                scoped.add(nn[0])
        regions = {pe.region_id for pe in plan.elements if pe.elem_id in eids}
        for pe in plan.elements:
            if pe.region_id in regions:
                scoped.add(pe.elem_id)
        return scoped

    def _region_metrics(self, plan, scoped: set) -> dict:
        """Local (region-scoped) harmony metrics — density vs page, alignment residual,
        gap rhythm, font-size tiers, pairwise contrast, local color ratio (T1/T2)."""
        from ppt_reflex.grid.composition import _color_ledger, _family_key, _rules
        from ppt_reflex.grid.color_utils import contrast_ratio

        elems = [e for e in plan.elements if e.elem_id in scoped]
        if not elems:
            return {}
        page_area = plan.page_w * plan.page_h
        page_density = round(sum(e.w * e.h for e in plan.elements) / page_area, 4) if page_area else 0.0
        xs = [e.x for e in elems]
        bbox_w = (max(e.x + e.w for e in elems) - min(xs)) or 1.0
        bbox_h = (max(e.y + e.h for e in elems) - min(e.y for e in elems)) or 1.0
        local_area = sum(e.w * e.h for e in elems)
        local_density = round(local_area / (bbox_w * bbox_h), 4)

        font_sizes = sorted({round(e.payload.font_size, 1) for e in elems
                             if e.payload and e.payload.font_size})
        left_edges = sorted(e.x for e in elems)
        align_residual = round(left_edges[-1] - left_edges[0], 1) if len(left_edges) > 1 else 0.0

        by_x = sorted(elems, key=lambda e: e.x)
        gaps = [by_x[i + 1].x - (by_x[i].x + by_x[i].w) for i in range(len(by_x) - 1)]
        gaps = [g for g in gaps if g >= 0]
        if len(gaps) >= 2:
            mean_g = sum(gaps) / len(gaps)
            gap_std = round((sum((g - mean_g) ** 2 for g in gaps) / len(gaps)) ** 0.5, 1)
        else:
            gap_std = 0.0

        pairwise = []
        for i in range(len(elems)):
            for j in range(i + 1, len(elems)):
                pi, pj = elems[i].payload, elems[j].payload
                if pi and pj and pi.font_color and pj.font_color:
                    pairwise.append({"a": elems[i].elem_id, "b": elems[j].elem_id,
                                     "ratio": round(contrast_ratio(pi.font_color, pj.font_color), 2)})

        ledger = [x for x in _color_ledger(plan, self._composition_context())
                  if x["elem_id"] in scoped]
        nt = _rules()["hue_harmony"]["neutral_chroma_threshold"]
        fam_areas: dict[str, float] = {}
        for x in ledger:
            key = _family_key(x["lch"], nt)
            fam_areas[key] = fam_areas.get(key, 0.0) + x["area"]
        total = sum(fam_areas.values())
        local_color_ratio = ({k: round(v / total, 3)
                              for k, v in sorted(fam_areas.items(), key=lambda kv: -kv[1])}
                             if total else {})

        return {
            "local_density": local_density,
            "page_density": page_density,
            "alignment_residual_pt": align_residual,
            "gap_sequence_std_pt": gap_std,
            "font_size_levels": font_sizes,
            "font_size_level_count": len(font_sizes),
            "pairwise_contrast": pairwise,
            "local_color_ratio": local_color_ratio,
        }

    def build_stream(self, path: str|None = None):
        """分页流式构建 — 逐页 yield 诊断，AI 可以立即看到每页结果。

        用法：
            for slide_result in b.build_stream("out.pptx"):
                ok = slide_result["ok"]
                errs = [d for d in slide_result["diagnostics"] if d["severity"] == "error"]
                print(f"S{slide_result['slide']:02d}: {len(errs)} errors")
                if errs:
                    break  # 中断构建，只修这一页

        与 build() 的区别：
        - build() 一次性完成所有页 + roundtrip，返回 dict
        - build_stream() 逐页 yield，不阻塞等待全部完成
        - 第一个返回的是 {'type': 'start', ...} 元信息
        - 之后每个是 {'type': 'slide', ...} 单页诊断
        - 最后是 {'type': 'summary', ...} 总汇 (含 roundtrip)
        """
        from pptx import Presentation; from pptx.util import Pt
        if path is None:
            ts = int(time.time()); path = os.path.join(tempfile.gettempdir(), f"ppt_reflex_{ts}.pptx")

        if self._template_pptx and os.path.exists(self._template_pptx):
            prs = Presentation(self._template_pptx)
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                prs.part.drop_rel(rId)
                prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
        else:
            prs = Presentation()
        prs.slide_width = Pt(self.pw); prs.slide_height = Pt(self.ph)
        total_slides = len(self._slides)

        yield {"type": "start", "total_slides": total_slides, "template": self._t.id, "style": self._style_id}

        all_diags: list[dict] = []
        for i, spec in enumerate(self._slides):
            slide_result = self.build_single_slide(i, prs=prs)
            all_diags.extend(slide_result.get("diagnostics", []))
            yield {**slide_result, "type": "slide", "slide": i, "slide_total": total_slides}

        prs.save(path)

        # P0-①: aggregate noisy diagnostics before returning
        aggregated_diags, agg_stats = _aggregate_diagnostics(all_diags)
        errs = [d for d in aggregated_diags if d.get("severity") in ("error",)]

        rt_results = check_overflow(path)
        for rt in rt_results:
            aggregated_diags.append(_diag(rt["slide"], "rt", None,
                                   kind=rt["kind"], severity=rt["severity"],
                                   message=rt["message"]))
        rt_errors = [rt for rt in rt_results if rt["severity"] == "error"]
        rt_warns = [rt for rt in rt_results if rt["severity"] == "warning"]

        diff_report = self.diff_log.diff()
        scope = self.diff_log.scope_alert()
        yield {
            "type": "summary",
            "path": path, "ok": len(errs) == 0 and len(rt_errors) == 0,
            "geometry_ok": len(errs) == 0 and len(rt_errors) == 0,
            "harmony_ok": _harmony_ok(all_diags),
            "diagnostics": aggregated_diags,
            "raw_diagnostic_count": agg_stats["raw_count"],
            "collapsed": {"dedup": agg_stats["dedup"],
                          "batch": agg_stats["batch"],
                          "trimmed_warnings": agg_stats["trimmed_warnings"],
                          "trimmed_info": agg_stats["trimmed_info"]},
            "summary": f"{agg_stats['final_count']} issues ({agg_stats['errors']} errors, "
                       f"{agg_stats['warnings']} warnings) — "
                       f"({agg_stats['raw_count']} raw, "
                       f"-{agg_stats['dedup']} dedup, "
                       f"{agg_stats['batch']} batch-collapsed)",
            "template": self._t.id, "style": self._style_id,
            "diff": {
                "entries": len(diff_report.entries) if diff_report else 0,
                "changed_elem_ids": list(diff_report.changed_elem_ids) if diff_report else [],
            },
            "scope_alert": scope,
        }

    def build(self, path: str|None = None) -> dict:
        from pptx import Presentation; from pptx.util import Pt
        if path is None:
            ts = int(time.time()); path = os.path.join(tempfile.gettempdir(), f"ppt_reflex_{ts}.pptx")

        # Use external template PPTX as base if provided (inherits master/layouts)
        if self._template_pptx and os.path.exists(self._template_pptx):
            prs = Presentation(self._template_pptx)
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                prs.part.drop_rel(rId)
                prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
        else:
            prs = Presentation()
        prs.slide_width = Pt(self.pw); prs.slide_height = Pt(self.ph)
        all_diags: list[dict] = []
        total_slides = len(self._slides)
        layout_fingerprints: list[str] = []

        self.diff_log.roll()  # DiffLog：上次 build 的 after 滚动为本次的 before
        for i, spec in enumerate(self._slides):
            plan = self._plan(spec); c = GridCanvas(GridConfig()); c.checkpoint()

            diags: list[dict] = []

            # Phase 0.5: region boundary validation
            plan.validate(verbose=False)
            for d in plan.diagnostics:
                diags.append(_diag(i, "0.5", d))

            # Phase 1: information layer layout + 全部后续阶段（单源 _pipeline_stages）
            execute_phase1(plan, c)
            self._pipeline_stages(i, plan, c, diags)

            # Fix #2: Render with smart layout selection
            _render_slide(prs, c, self._t, slide_index=i, total_slides=total_slides,
                          slide_spec=spec)

            self.diff_log.snap_after(plan, i)

            all_diags.extend(diags)
            layout_fingerprints.append(_layout_fingerprint(plan))

        orig_path = path
        path = _ensure_writable_path(path)
        prs.save(path)
        if path != orig_path:
            print(f"[PPTBuilder] WARNING: '{orig_path}' is locked (open in PowerPoint/WPS) — wrote to '{path}'")
        errs = [d for d in all_diags if d.get("severity") in ("error",)]
        warns = [d for d in all_diags if d.get("severity") in ("warning","warn")]

        # Post-render roundtrip: reopen saved PPTX and verify every text box fits
        rt_results = check_overflow(path)
        for rt in rt_results:
            all_diags.append(_diag(rt["slide"], "rt", None,
                                   kind=rt["kind"], severity=rt["severity"],
                                   message=rt["message"]))
        rt_errors = [rt for rt in rt_results if rt["severity"] == "error"]
        rt_warns = [rt for rt in rt_results if rt["severity"] == "warning"]

        diff_report = self.diff_log.diff()
        scope = self.diff_log.scope_alert()

        # P0-①: aggregate noisy diagnostics before returning
        aggregated_diags, agg_stats = _aggregate_diagnostics(all_diags)
        errs = [d for d in aggregated_diags if d.get("severity") in ("error",)]

        # P0+P1: design policy + circuit breaker — detects fix-loops and systemic issues
        from ppt_reflex.design_policy import analyze_design_issues, gather_slides_data
        slides_data = gather_slides_data(self)
        design_hints = analyze_design_issues(aggregated_diags, all_diags, slides_data,
                                             breaker=self._breaker)

        # Hard block: if breaker has escalated to BLOCK, ok=False regardless
        blocked = self._breaker.blocked_fingerprints()
        hard_blocked = len(blocked) > 0 and self._breaker.build_count >= 3
        build_ok = len(errs) == 0 and len(rt_errors) == 0 and not hard_blocked

        # T8: human-panel mode (strict_tokens=False) may introduce WCAG contrast
        # violations — surface them as a non-blocking warning, never override the human.
        human_override_warning = (not self.strict_tokens) and _has_contrast_violation(all_diags)

        # T5: two-level correctness. geometry_ok = existing ok (zero geometric errors);
        # harmony_ok = zero error/warning among T2–T4 harmony rules. ok stays geometry_ok.
        harmony_ok = _harmony_ok(all_diags)

        return {"path": path, "ok": build_ok,
                "geometry_ok": build_ok,
                "harmony_ok": harmony_ok,
                "human_override_warning": human_override_warning,
                "diagnostics": aggregated_diags,
                "design_hints": design_hints,
                "build_number": self._breaker.build_count,
                "hard_blocked": hard_blocked,
                "blocked_fingerprints": [b if isinstance(b, dict) else b._asdict() for b in blocked],
                "entropy_stalled": self._breaker.is_stalled,
                "raw_diagnostic_count": agg_stats["raw_count"],
                "collapsed": {"dedup": agg_stats["dedup"],
                              "batch": agg_stats["batch"],
                              "trimmed_warnings": agg_stats["trimmed_warnings"],
                              "trimmed_info": agg_stats["trimmed_info"]},
                "summary": f"{agg_stats['final_count']} issues ({agg_stats['errors']} errors, "
                           f"{agg_stats['warnings']} warnings) — "
                           f"({agg_stats['raw_count']} raw, "
                           f"-{agg_stats['dedup']} dedup, "
                           f"{agg_stats['batch']} batch-collapsed, "
                           f"-{agg_stats['trimmed_warnings']} warn / -{agg_stats['trimmed_info']} info trimmed)",
                "template": self._t.id, "style": self._style_id,
                "diff": {
                    "entries": len(diff_report.entries) if diff_report else 0,
                    "changed_elem_ids": list(diff_report.changed_elem_ids) if diff_report else [],
                },
                "scope_alert": scope}

    def clear_diff_log(self):
        """User confirmed deck is done — wipe mutation trace."""
        self.diff_log.clear()

    # ── P0: verify() — post-render structural verification ──

    def verify(self, path: str) -> dict:
        """Reopen a generated PPTX and verify structural integrity without needing to
        view it. Returns per-slide metrics an AI can act on:
          - decoration presence (frame/rail/bar/dot actually drawn)
          - element counts & overlap
          - content coverage (whitespace proxy)
          - page number sanity
        No visual model needed — pure geometry from the saved file.
        """
        from pptx import Presentation
        from pptx.util import Emu

        if not os.path.exists(path):
            return {"ok": False, "error": f"file not found: {path}", "slides": []}

        prs = Presentation(path)
        page_w = prs.slide_width / Emu(1)
        page_h = prs.slide_height / Emu(1)
        slides_out = []

        for i, slide in enumerate(prs.slides):
            shapes = list(slide.shapes)
            full_w = prs.slide_width
            borders = [s for s in shapes if s.width == full_w
                       and s.height < Emu(400000)]  # full-width thin band
            has_top = any(s.top < Emu(200000) for s in borders)
            has_bottom = any(s.top > prs.slide_height - Emu(400000) for s in borders)

            texts = [s for s in shapes if s.has_text_frame]
            text_chars = sum(len(s.text_frame.text) for s in texts)

            # coverage: non-background shapes (exclude full-width borders which are bg-like)
            content = [s for s in shapes if not (s.width == full_w and s.height < Emu(400000))]
            cov = min(1.0, sum((s.width / Emu(1)) * (s.height / Emu(1))
                               for s in content) / (page_w * page_h))

            # ── 形状双向检测 ──
            # 正向: 形状内文字是否水平+垂直居中 (shape → text)
            # 反向: 自由 textbox 中心是否落在某形状内 (text → shape 归属)
            from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
            shape_cells = []      # 形状内文字验证结果
            orphan_texts = []     # 反向: 不在任何形状内的文字

            for s in shapes:
                if not s.has_text_frame:
                    continue
                t = s.text_frame.text.strip()
                if not t:
                    continue
                # 判断是否带 fill 的形状 (非纯 textbox)
                has_fill = False
                try:
                    from pptx.enum.dml import MSO_FILL_TYPE
                    ft = s.fill.type
                    if ft is not None and ft not in (MSO_FILL_TYPE.BACKGROUND, MSO_FILL_TYPE.GROUP):
                        has_fill = True
                except Exception:
                    has_fill = False
                if not has_fill:
                    # 纯 textbox — 反向检测: 中心是否落在某形状内
                    cx = (s.left + s.width / 2) / Emu(1)
                    cy = (s.top + s.height / 2) / Emu(1)
                    inside = False
                    for other in shapes:
                        if other is s or not other.has_text_frame:
                            continue
                        try:
                            if other.fill.type is None:
                                continue
                        except Exception:
                            continue
                        if (other.left <= s.left + s.width / 2 <= other.left + other.width and
                                other.top <= s.top + s.height / 2 <= other.top + other.height):
                            inside = True
                            break
                    if not inside:
                        orphan_texts.append({
                            "text": t[:20], "x": round(cx), "y": round(cy),
                        })
                    continue

                # 带 fill 的形状 — 正向检测居中
                # 跳过 rounded_rectangle：b.box() 卡片是文本容器，左对齐正常；
                # 只有语义形状（hexagon/diamond/oval/chevron…）才要求形状内文字居中。
                try:
                    from pptx.enum.shapes import MSO_SHAPE
                    if getattr(s, "auto_shape_type", None) == MSO_SHAPE.ROUNDED_RECTANGLE:
                        continue
                except Exception:
                    pass
                try:
                    align_h = s.text_frame.paragraphs[0].alignment
                    v_anchor = s.text_frame.vertical_anchor
                    h_ok = align_h == PP_ALIGN.CENTER
                    v_ok = v_anchor == MSO_ANCHOR.MIDDLE
                    if h_ok and v_ok:
                        centering = "centered"
                    elif h_ok:
                        centering = "v_not_centered"
                    elif v_ok:
                        centering = "h_not_centered"
                    else:
                        centering = "not_centered"
                except Exception:
                    centering = "unknown"
                shape_cells.append({
                    "text": t[:20], "centering": centering,
                })

            slides_out.append({
                "slide": i + 1,
                "n_shapes": len(shapes),
                "n_text": len(texts),
                "n_borders": len(borders),
                "frame_top": has_top,
                "frame_bottom": has_bottom,
                "coverage": round(cov, 3),
                "text_chars": text_chars,
                "shape_cells": shape_cells,
                "orphan_texts": orphan_texts,
            })

        return {"ok": True, "path": path, "page_w": page_w, "page_h": page_h,
                "slides": slides_out}

    # ── P1-①: incremental rebuild ──

    def rebuild(self, changed_slides: list[int], path: str|None = None) -> dict:
        """Only re-run pipeline for the given slide indices; others reuse pipelined plan+canvas.

        Call this after fix_slide() — much faster than build() for 1-2 slide edits.
        Uses _pipeline_cache which stores (hash, LayoutPlan, GridCanvas, diags) per slide.
        """
        from pptx import Presentation
        from pptx.util import Pt

        if path is None:
            ts = int(time.time())
            path = os.path.join(tempfile.gettempdir(), f"ppt_reflex_{ts}.pptx")

        changed: set[int] = set(changed_slides)
        total_slides = len(self._slides)
        all_diags: list[dict] = []

        if self._template_pptx and os.path.exists(self._template_pptx):
            prs = Presentation(self._template_pptx)
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                prs.part.drop_rel(rId)
                prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
        else:
            prs = Presentation()
        prs.slide_width = Pt(self.pw)
        prs.slide_height = Pt(self.ph)

        self.diff_log.roll()
        for i, spec in enumerate(self._slides):
            h = self._spec_hash(spec)

            if i not in changed:
                cached_h, cached_plan, cached_canvas, cached_diags = self._pipeline_cache.get(i, (None, None, None, None))
                if cached_h == h and cached_plan is not None and cached_canvas is not None:
                    # Cache hit: reuse pipeline results, just re-render
                    _render_slide(prs, cached_canvas, self._t, slide_index=i, total_slides=total_slides,
                                  slide_spec=spec)
                    all_diags.extend(cached_diags or [])
                    # Update cache with fresh slide XML
                    self._pipeline_cache[i] = (h, cached_plan, cached_canvas, cached_diags)
                    continue

            # Cache miss or changed slide → run full pipeline
            plan = self._plan(spec)
            c = GridCanvas(GridConfig())
            c.checkpoint()

            diags: list[dict] = []

            plan.validate(verbose=False)
            for d in plan.diagnostics:
                diags.append(_diag(i, "0.5", d))

            execute_phase1(plan, c)
            self._pipeline_stages(i, plan, c, diags)

            _render_slide(prs, c, self._t, slide_index=i, total_slides=total_slides,
                          slide_spec=spec)
            self.diff_log.snap_after(plan, i)

            # Cache pipeline results
            self._pipeline_cache[i] = (h, plan, c, diags)
            all_diags.extend(diags)

        prs.save(path)
        rt_results = check_overflow(path)
        for rt in rt_results:
            all_diags.append(_diag(rt["slide"], "rt", None,
                                   kind=rt["kind"], severity=rt["severity"],
                                   message=rt["message"]))
        rt_errors = [rt for rt in rt_results if rt["severity"] == "error"]

        diff_report = self.diff_log.diff()
        scope = self.diff_log.scope_alert()

        aggregated_diags, agg_stats = _aggregate_diagnostics(all_diags)
        errs = [d for d in aggregated_diags if d.get("severity") in ("error",)]

        # P0+P1: design policy + circuit breaker
        from ppt_reflex.design_policy import analyze_design_issues, gather_slides_data
        slides_data = gather_slides_data(self)
        design_hints = analyze_design_issues(aggregated_diags, all_diags, slides_data,
                                             breaker=self._breaker)
        blocked = self._breaker.blocked_fingerprints()
        hard_blocked = len(blocked) > 0 and self._breaker.build_count >= 3
        build_ok = len(errs) == 0 and len(rt_errors) == 0 and not hard_blocked

        # T8: human-panel mode contrast violations — non-blocking warning.
        human_override_warning = (not self.strict_tokens) and _has_contrast_violation(all_diags)

        # T5: two-level correctness.
        harmony_ok = _harmony_ok(all_diags)

        return {
            "path": path,
            "ok": build_ok,
            "geometry_ok": build_ok,
            "harmony_ok": harmony_ok,
            "human_override_warning": human_override_warning,
            "diagnostics": aggregated_diags,
            "design_hints": design_hints,
            "build_number": self._breaker.build_count,
            "hard_blocked": hard_blocked,
            "blocked_fingerprints": [b if isinstance(b, dict) else b._asdict() for b in blocked],
            "entropy_stalled": self._breaker.is_stalled,
            "raw_diagnostic_count": agg_stats["raw_count"],
            "collapsed": {"dedup": agg_stats["dedup"],
                          "batch": agg_stats["batch"],
                          "trimmed_warnings": agg_stats["trimmed_warnings"],
                          "trimmed_info": agg_stats["trimmed_info"]},
            "summary": f"{agg_stats['final_count']} issues ({agg_stats['errors']} errors) — "
                       f"{len(changed)} slides rebuilt, "
                       f"{total_slides - len(changed)} from cache",
            "template": self._t.id,
            "style": self._style_id,
            "diff": {
                "entries": len(diff_report.entries) if diff_report else 0,
                "changed_elem_ids": list(diff_report.changed_elem_ids) if diff_report else [],
            },
            "scope_alert": scope,
        }

    def clear_cache(self):
        """Drop incremental rebuild cache (e.g. after template switch)."""
        self._pipeline_cache.clear()

    def clear_circuit_breaker(self):
        """Reset fix-loop tracking — call after a DESIGN-LEVEL change (template switch,
        layout redesign, content rewrite). Does NOT clear pipeline cache."""
        self._breaker.reset()

    def declare_direction(self, direction: str) -> str | None:
        """Declare the fix direction for the upcoming build()/rebuild() call.

        Valid directions: increase_box_height, increase_box_width, decrease_font_size,
        increase_region, rearrange_regions, reduce_text, split_text, shorter_lines,
        remove_elements, split_slide, switch_layout, switch_region_order,
        change_text_color, change_fill_color, switch_template, switch_style,
        dark_to_light, light_to_dark.

        Returns error string if invalid, None if OK.
        """
        return self._breaker.declare_direction(direction)


    def title(self, text: str, region: str = "main") -> _Spec:
        return self._s("Heading", text, region, "text", ph=40)
    def subtitle(self, text: str, region: str = "main") -> _Spec:
        return self._s("Subtitle", text, region, "text", ph=30)
    def text(self, text: str, style: str = "Body", region: str = "main") -> _Spec:
        return self._s(style, text, region, "text")
    def bullet(self, text: str, region: str = "main") -> _Spec:
        return self._s("ListItem", f"• {text}", region, "text")
    def footer(self, text: str, region: str = "footer") -> _Spec:
        return self._s("Footer", text, region, "footer")
    def box(self, text: str, style: str = "Body", region: str = "main",
            fill_color: tuple|None = None, shape_id: str|None = None,
            ph: float|None = None, align_h: str = "left", allow_shrink: bool = False,
            role: str = "", recipe: str|None = None,
            corner_radius: float|None = None, **kwargs) -> _Spec:
        # T7: box(radius=) → corner_radius (CSS alias); unknown kwargs rejected.
        radius = kwargs.pop("radius", None)
        if radius is not None:
            corner_radius = radius
        reject_unknown_kwargs("box", kwargs)
        # T8: agent mode rejects element-level raw color; recipe/style tiers own color.
        if self.strict_tokens and fill_color is not None:
            _forbid_raw_color(
                "box(fill_color=...)",
                "Use recipe='card'/'kpi'/'quote' (or a style tier) to carry color.",
            )
        # recipe: named component from recipes.json (card/kpi/quote) — token values
        # are pre-resolved by the design-token layer; explicit args still win.
        if recipe:
            from ppt_reflex.grid.design_tokens import resolve_recipe
            rec = resolve_recipe(recipe)
            if shape_id is None and rec.get("shape"):
                shape_id = rec["shape"]
            if fill_color is None:
                c = rec.get("fill")
                if isinstance(c, str) and c.startswith("#"):
                    fill_color = _hex_to_rgb(c)
            if corner_radius is None and isinstance(rec.get("radius"), (int, float)):
                corner_radius = rec["radius"]
            if align_h == "left" and rec.get("align"):
                align_h = str(rec["align"]).lower()
        if shape_id is None:
            shape_id = "rounded_rectangle"
        return self._s(style, text, region, "textbox", fill_color=fill_color,
                       shape_id=shape_id, ph=ph, align_h=align_h, allow_shrink=allow_shrink,
                       role=role, corner_radius=corner_radius)
    def shape(self, shape_id: str, region: str = "main",
              fill_color: tuple|None = None, pw: float|None = None, ph: float|None = None,
              text: str = "", font_size: float|None = None,
              font_color: tuple|None = None, align_h: str = "center",
              role: str = "", **kwargs) -> _Spec:
        # T7: strict mode rejects CSS hallucinations on shape().
        reject_unknown_kwargs("shape", kwargs)
        # T8: agent mode rejects raw color and the type-scale backdoor.
        if self.strict_tokens:
            if fill_color is not None:
                _forbid_raw_color("shape(fill_color=...)", "Use recipe/style tiers for color.")
            if font_color is not None:
                _forbid_raw_color("shape(font_color=...)",
                    "font_color_override bypasses the style/type-scale contract — use a style tier.")
            if font_size is not None:
                _forbid_raw_color("shape(font_size=...)",
                    "font_size_override bypasses the type scale — use a style tier (Heading/Body/…).")
        # 形状可承载文字（圆形数字/步骤节点/品牌标）。文字居中在形状内，无需额外文本框。
        # text 非空 + 无 fill → 透明底形状只做文字容器（文字可自由摆放在形状上）。
        style = "Body"
        if text and font_size is not None:
            style = "__shape_inline__"
        elif text:
            style = "Emphasis"
        return _Spec(elem_id=self._nid("shape"), style=style, text=text, region=region,
                     ctype="shape", fill_color=fill_color, shape_id=shape_id, pw=pw, ph=ph,
                     align_h=align_h, font_size_override=font_size,
                     font_color_override=font_color, role=role)
    def image(self, path: str, region: str = "main",
              pw: float|None = None, ph: float|None = None,
              fit_mode: str = "fit", allow_upscale: bool = False,
              layout_mode: str = "", caption: str = "", **kwargs) -> _Spec:
        # T7: strict mode rejects CSS hallucinations; "contain"/"cover" → fit/fill.
        reject_unknown_kwargs("image", kwargs)
        fit_mode = normalize_fit_mode(fit_mode)
        # Fix #10: validate path exists
        if not os.path.isfile(path):
            print(f"[PPTBuilder] WARNING: image path not found: {path}")
        return _Spec(elem_id=self._nid("img"), style="", text="", region=region,
                     ctype="image", pw=pw, ph=ph, image_path=path,
                     fit_mode=fit_mode, allow_upscale=allow_upscale,
                     layout_mode=layout_mode, caption=caption)
    def table(self, headers: list[str], rows: list[list[str]],
              region: str = "main", font_size: float = 12.0,
              header_bg: tuple|None = None) -> _Spec:
        """Add a table element. headers=column names, rows=row data.

        Each row must have len(headers) cells. The engine auto-sizes columns
        to fit region width and auto-sizes rows by content.
        """
        return _Spec(elem_id=self._nid("tbl"), style="Table", text="", region=region,
                     ctype="table", ph=None,
                     table_headers=list(headers), table_rows=[list(r) for r in rows])
    def arrow(self, frm: str, to: str, text: str = "", direction: str = "below",
              color: tuple = None, width: float = 1.5,
              margin_pt: float = 8.0, text_font_size: float = 10.0,
              text_color: tuple = None,
              occlusion_check: bool = True) -> _Arrow:
        # T8: explicit color = agent raw-color token → rejected; None = template contract.
        if self.strict_tokens and (color is not None or text_color is not None):
            _forbid_raw_color(
                "arrow(color=.../text_color=...)",
                "The arrow color is the template contract — omit color/text_color to inherit it.",
            )
        # Allow passing _Spec objects directly — resolve to elem_id
        from_eid = frm.elem_id if hasattr(frm, 'elem_id') else frm
        to_eid = to.elem_id if hasattr(to, 'elem_id') else to
        return _Arrow(self._nid("arrow"), from_eid, to_eid, text, direction, color, width,
                      margin_pt, text_font_size, text_color, occlusion_check)
    def divider(self, region: str = "main", color: tuple|None = None, width_pt: float|None = None) -> _Spec:
        # T8: agent mode rejects a raw divider color; the template owns divider_color_hex.
        if self.strict_tokens and color is not None:
            _forbid_raw_color(
                "divider(color=...)",
                "The divider color is the template contract (divider_color_hex) — omit color to inherit it.",
            )
        dh = self._t.divider_color_hex
        if dh:
            c = color or tuple(int(dh.lstrip("#")[i:i+2], 16) for i in (0, 2, 4))
        else:
            c = color or (0x44, 0x44, 0x55)
        # 线宽默认取模板契约（divider_width_pt），调用方可显式覆盖
        w_pt = width_pt if width_pt is not None else (self._t.divider_width_pt or 3.0)
        return _Spec(elem_id=self._nid("div"), style="", text="", region=region,
                     ctype="shape", fill_color=c, shape_id="rectangle", ph=w_pt)

    # ── Image layout auto-inference ──
    def auto_layout_mode(self, image_path: str) -> str:
        """Select layout mode from image aspect ratio + preset constraints. Decision tree:
        aspect > 1.6 -> hero_top (landscape)
        aspect < 0.8 -> hero_right (portrait)
        aspect 0.8-1.6 -> center_float (square)
        Falls back to preset's preferred_modes[0] if decision mode not in allowed set."""
        from PIL import Image
        try:
            img = Image.open(image_path)
            w, h = img.size
            aspect = w / h if h > 0 else 1.0
        except Exception as e:
            import traceback
            traceback.print_exc()
            print(f"[PPTBuilder] auto_layout_mode PIL error for {image_path}: {e}")
            return "center_float"

        # Decide by aspect ratio
        if aspect > 1.6:
            mode = "hero_top"
        elif aspect < 0.8:
            mode = "hero_right"
        else:
            mode = "center_float"

        # If preset constrains preferred_modes and decision mode is not in it, fall back to preferred_modes[0]
        if self._image_layout:
            preferred = self._image_layout.get("preferred_modes", [])
            if preferred and mode not in preferred:
                mode = preferred[0]
        return mode

    def image_constraints(self, layout_mode: str) -> dict:
        """Return max_w/max_h/anchor/ratio constraints for the given layout_mode from current preset."""
        mc = self._image_layout.get("mode_constraints", {}) if self._image_layout else {}
        return mc.get(layout_mode, {})

    def image_treatment(self) -> dict:
        """Return image rendering treatment (corner_radius/border/shadow) from current preset."""
        defaults = {"corner_radius_pt": 0, "border_role": "none", "shadow_role": "none"}
        if self._image_layout:
            defaults.update(self._image_layout.get("treatment", {}))
        return defaults

    def caption_format(self) -> dict:
        """Return caption format convention from current preset."""
        defaults = {"font_size": 11, "alignment": "left", "max_lines": 1, "prefix": ""}
        if self._image_layout:
            defaults.update(self._image_layout.get("caption", {}))
        return defaults

    # ── Internals ──
    def _nid(self, pfx: str) -> str:
        self._id += 1; return f"{pfx}_{self._id}"

    def _s(self, style: str, text: str, region: str, ctype: str,
           fill_color=None, shape_id="", ph=None,
           align_h="left", allow_shrink=False, allow_wrap=False, role: str = "",
           corner_radius: float|None = None) -> _Spec:
        s = STYLE.get(style, STYLE["Body"])
        return _Spec(self._nid("e"), style, text, region, ctype, "stack",
                     fill_color=fill_color or s.get("fill_color"),
                     shape_id=shape_id, ph=ph, margin=4.0,
                     align_h=align_h, allow_shrink=allow_shrink, allow_wrap=allow_wrap,
                     role=role, corner_radius=corner_radius)

    # Fix #1: resolve style colors + font sizes from template profile
    def _resolve_style(self, style_name: str, fill_color=None) -> dict:
        """Merge STYLE defaults with template/style-preset colors AND font sizes.
        优先级：template/style preset > STYLE 表。"""
        s = dict(STYLE.get(style_name, STYLE["Body"]))
        t = self._t

        # ── ① Template / style preset 层 ──
        # Map style -> template color field
        if style_name in ("Heading", "Subheading"):
            s["font_color"] = _hex_to_rgb(t.title_hex)
        elif style_name == "Subtitle":
            s["font_color"] = _hex_to_rgb(t.gray_hex) if t.gray_hex else _hex_to_rgb(t.text_hex)
        elif style_name == "Emphasis":
            s["font_color"] = _hex_to_rgb(t.accent2_hex) if t.accent2_hex else _hex_to_rgb(t.accent_hex)
        elif style_name == "Caption":
            s["font_color"] = _hex_to_rgb(t.gray_hex) if t.gray_hex else _hex_to_rgb(t.text_hex)
        elif style_name == "Footer":
            s["font_color"] = _hex_to_rgb(t.dim_hex) if t.dim_hex else _hex_to_rgb(t.gray_hex)
        else:  # "Body", "ListItem"
            s["font_color"] = _hex_to_rgb(t.text_hex)

        # Map style -> template font size
        # Subtitle 用独立 subtitle_size（0=auto→body+4）；旧版错接到 caption_size，
        # 导致副标题比正文还小（academic: subtitle 14 < body 20）
        if style_name in ("Heading",):
            s["font_size"] = t.title_size
        elif style_name in ("Subtitle",):
            s["font_size"] = t.subtitle_size or (t.body_size + 4)
        elif style_name in ("Subheading",):
            s["font_size"] = t.body_size + 2
        elif style_name in ("Body", "ListItem", "Emphasis"):
            s["font_size"] = t.body_size
        elif style_name == "Caption":
            s["font_size"] = t.caption_size
        elif style_name == "Footer":
            s["font_size"] = t.page_number_size

        # center_titles 模板契约：声明不居中的模板，标题左对齐（旧版永远居中）
        if style_name == "Heading" and not t.center_titles:
            s["alignment"] = "LEFT"

        # Fix #13: font_name fallback
        if "font_name" not in s or not s["font_name"]:
            s["font_name"] = self._style_body_font or t.body_font

        # Dark fill -> white text
        if fill_color and _is_dark(fill_color):
            s["font_color"] = (0xFF, 0xFF, 0xFF)

        return s

    def _plan(self, spec: _Slide) -> LayoutPlan:
        ctmap = {"text": ContentType.TEXT, "textbox": ContentType.TEXTBOX,
                 "shape": ContentType.SHAPE, "image": ContentType.IMAGE,
                 "table": ContentType.TABLE,
                 "annotation": ContentType.ANNOTATION, "footer": ContentType.FOOTER}
        regions = []
        for ri, d in enumerate(spec.regions):
            rid, x, y, w, h = d[0], d[1], d[2], d[3], d[4]
            ro = d[5] if len(d) > 5 else ri + 1
            # Fix #5: optional 6th field for content_inset
            inset = d[6] if len(d) > 6 else 8.0
            regions.append(Region(rid, x, y, w, h, rid, ro, inset))

        elems = []
        for e in spec.elements:
            # Fix #1: resolve style from template colors
            s = self._resolve_style(e.style, e.fill_color)
            # shape() 文字 override — 压过解析值（font_size_override / font_color_override）
            if e.font_size_override is not None:
                s["font_size"] = e.font_size_override
            if e.font_color_override is not None:
                s["font_color"] = e.font_color_override
            # 显式 align_h 覆盖 style 默认对齐 — shape(text=...) 默认 center，确保形状内文字水平居中
            if e.align_h and e.align_h != "left":
                s["alignment"] = e.align_h.upper()
            ctyp = ctmap.get(e.ctype, ContentType.TEXT)

            # ── layout_mode → 图片尺寸/锚点约束（此前 7 种布局模式渲染无差异）──
            pw, ph, align_h = e.pw, e.ph, e.align_h
            if ctyp == ContentType.IMAGE and e.layout_mode:
                mc = self.image_constraints(e.layout_mode)
                if mc:
                    reg = next((r for r in regions if r.region_id == e.region), None)
                    if reg is not None:
                        _, _, rw, rh = reg.usable_rect
                        if pw is None and mc.get("max_width_pt"):
                            pw = min(rw, float(mc["max_width_pt"]))
                        if ph is None and mc.get("max_height_pt"):
                            ph = min(rh, float(mc["max_height_pt"]))
                    if mc.get("anchor") in ("center", "top_center") and align_h == "left":
                        align_h = "center"

            # ── 圆角半径：显式 corner_radius → style preset shape_override → template.card_rounding ──
            corner_radius_pt = e.corner_radius
            if corner_radius_pt is None and e.shape_id == "rounded_rectangle":
                cr = (self._shape_override.get("corner_radius_pt") or {})
                corner_radius_pt = cr.get("card") or (self._t.card_rounding or None)

            # ── 语义 role 接线（"AI 写语义"的落点）：AI 可声明 emphasis/backdrop/
            # connector 等 overlay 角色，碰撞系统据此放行；未声明 → 族扶手默认
            declared_role = None
            if e.role:
                try:
                    declared_role = SemanticRole(e.role.lower())
                except ValueError:
                    print(f"[PPTBuilder] WARNING: unknown role '{e.role}' on {e.elem_id}, "
                          f"falling back to family default")

            p = ElementPayload(
                role=declared_role,
                text=e.text,
                style_name=e.style if e.style != "__shape_inline__" else "",
                font_size=s.get("font_size", 14),
                font_color=s.get("font_color", (0x33, 0x33, 0x44)),
                font_bold=s.get("font_bold", False),
                font_name=s.get("font_name", self._t.body_font),
                alignment=s.get("alignment", "LEFT"),
                fill_color=e.fill_color or s.get("fill_color"),
                shape_id=e.shape_id,
                corner_radius_pt=corner_radius_pt,
                line_spacing=self._t.line_spacing,
                image_path=e.image_path,
                fit_mode=e.fit_mode,
                allow_upscale=e.allow_upscale,
                layout_mode=e.layout_mode,
                caption=e.caption,
                table_headers=e.table_headers if e.table_headers else None,
                table_rows=e.table_rows if e.table_rows else None,
            )
            # Fix #6: expose Phase1Element params
            elems.append(Phase1Element(
                e.elem_id, e.region, ctyp,
                payload=p, fill_mode=e.fill_mode, margin_above=e.margin,
                preferred_width=pw, preferred_height=ph,
                align_h=align_h,
                allow_shrink=e.allow_shrink,
                allow_wrap=e.allow_wrap,
                ARROW_SLOT=e.arrow_slot,
            ))

        decos = []
        for a in spec.arrows:
            # Fix #8: pass all DecoIntent params
            # T8: None color → template contract color (engine owns the hue).
            _arrow_line = a.color if a.color is not None else _hex_to_rgb(
                self._t.divider_color_hex or self._t.accent_hex or "666666")
            _arrow_text = a.text_color if a.text_color is not None else _hex_to_rgb(
                self._t.gray_hex or self._t.text_hex or "555555")
            decos.append(DecoIntent(
                a.deco_id, "arrow", [a.from_elem, a.to_elem],
                a.direction,
                margin_pt=a.margin_pt,
                style={"line_color": _arrow_line, "line_width_pt": a.width},
                text=a.text,
                text_font_size=a.text_font_size,
                text_color=_arrow_text,
                occlusion_check=a.occlusion_check,
            ))

        # Fix #7: use template page_margin for page_safe_inset
        plan = LayoutPlan(
            page_w=self.pw, page_h=self.ph,
            page_safe_inset=max(12.0, self._t.page_margin),
            title=spec.title,
            regions=regions,
            phase1_elements=elems,
            deco_intents=decos,
        )
        return plan

    def _composition_context(self) -> dict:
        """Resolved template+style palette for the harmony floor (T2–T4)."""
        return {
            "bg_hex": self._t.bg_hex,
            "accent_hex": self._t.accent_hex,
            "accent2_hex": self._t.accent2_hex,
            "text_hex": self._t.text_hex,
            "gray_hex": self._t.gray_hex,
            "style_id": self._style_id,
        }

    def _run_aesthetics(self, canvas, plan) -> list[dict]:
        """Run AestheticsEngine and return structured diagnostics."""
        engine = AestheticsEngine()
        elems = []
        for pe in plan.elements:
            p = pe.payload
            fill_hex = _rgb_to_hex(p.fill_color) if p and p.fill_color else self._t.bg_hex.lstrip("#")
            font_hex = _rgb_to_hex(p.font_color) if p and p.font_color else "000000"
            es = ElemStyle(
                id=pe.elem_id, content_type=pe.content_type,
                font_size=p.font_size if p else 12,
                font_bold=p.font_bold if p else False,
                font_color=font_hex, fill_color=fill_hex,
                line_spacing=p.line_spacing if p else 1.2,
                text=p.text if p else "",
                x=pe.x, y=pe.y, w=pe.w, h=pe.h,
                auto_size="SHAPE_TO_FIT_TEXT",
                canvas_w=self.pw, canvas_h=self.ph,
            )
            # No fill -> inherit slide background (avoids false contrast errors when text over slide bg)
            if p and not p.fill_color and not p.shape_id:
                es.fill_color = self._t.bg_hex
            if p and p.fill_color:
                es.fill_color = _rgb_to_hex(p.fill_color)
            elems.append(es)

        violations = engine.check(elems, timing="commit",
                                  ctx={"max_elements": self._t.max_elements_per_slide,
                                       "max_chars": self._t.max_chars_per_slide,
                                       "max_colors": self._t.max_colors})
        return [_ae_violation_to_diag(v) for v in violations]


def _layout_fingerprint(plan) -> str:
    """Layout signature of one slide — region topology, order-insensitive.

    Buckets each region by its grid position (top/bottom × left/center/right) + size class.
    Two slides with the same bucket multiset are considered "same layout". """
    buckets = []
    for r in plan.sorted_regions():
        cy = r.y + r.h / 2
        cx = r.x + r.w / 2
        v = "top" if cy < plan.page_h / 3 else ("bottom" if cy > 2 * plan.page_h / 3 else "mid")
        h = "left" if cx < plan.page_w / 3 else ("right" if cx > 2 * plan.page_w / 3 else "center")
        size = "big" if r.w * r.h > (plan.page_w * plan.page_h) / 4 else "small"
        buckets.append(f"{v}-{h}-{size}")
    return "|".join(sorted(buckets))


def _activate_collision(plan, canvas, slide_idx: int) -> list[dict]:
    """Activate the dormant three-layer collision system.

    After phase1 locks pt coordinates, occupy every element's bbox into the
    info_grid, then pairwise-check overlaps. Severity comes from POLICIES
    (TEXT×TEXT → error; BAND×BAND → warn; overlay families → allow). Overlaps
    between elements of DIFFERENT regions are flagged as cross-region (the
    case that mattered all along: content bleeding out of its region).

    Returns diagnostics. Occupies info_grid (feeds supply/spatial/profiles) but
    never mutates the locked layout coordinates.
    """
    from ppt_reflex.grid.types import _verdict_to_level
    from ppt_reflex.grid.plan import PageElement

    rects = getattr(canvas, "_phase1_rects", {})
    if not rects:
        return []

    diags: list[dict] = []

    # Occupy every element into the info_grid — feeds the dormant three-layer
    # system (supply/spatial/profiles read occupied state) AND enables collision.
    from ppt_reflex.grid.types import SemanticRole
    for pe in plan.elements:
        if not isinstance(pe, PageElement):
            continue
        role = SemanticRole.ENTITY
        try:
            pol = POLICIES[family_of(pe.content_type)]
            role = pol.default_role
        except KeyError:
            pass
        # AI 声明的语义 role 优先于族扶手（"AI 写语义"——此前一律被族默认覆盖）
        declared = getattr(getattr(pe, "payload", None), "role", None)
        if declared is not None:
            role = declared
        canvas.info_grid.occupy_bbox(
            pe.x, pe.y, pe.w, pe.h, pe.elem_id, pe.content_type,
            role=role, source="engine",
        )

    # region_id per element (for cross-region flagging)
    elem_region = {}
    for pe in plan.elements:
        if isinstance(pe, PageElement):
            elem_region[pe.elem_id] = getattr(pe, "region_id", "")

    elems = list(plan.elements)
    for i in range(len(elems)):
        for j in range(i + 1, len(elems)):
            a, b = elems[i], elems[j]
            if not isinstance(a, PageElement) or not isinstance(b, PageElement):
                continue
            ax, ay, aw, ah = a.x, a.y, a.w, a.h
            bx, by, bw, bh = b.x, b.y, b.w, b.h
            if ax >= bx + bw or bx >= ax + aw or ay >= by + bh or by >= ay + ah:
                continue

            fam_a = family_of(a.content_type)
            fam_b = family_of(b.content_type)

            # Overlay families never collide (CONNECTOR/EMPHASIS/BACKDROP)
            if fam_a in (Family.CONNECTOR, Family.EMPHASIS, Family.BACKDROP) or \
               fam_b in (Family.CONNECTOR, Family.EMPHASIS, Family.BACKDROP):
                continue

            # AI 声明 overlay role 的元素不参与碰撞（emphasis 高亮框、backdrop 衬底…）
            role_a = getattr(getattr(a, "payload", None), "role", None)
            role_b = getattr(getattr(b, "payload", None), "role", None)
            if (role_a is not None and role_a in OVERLAY_ROLES) or \
               (role_b is not None and role_b in OVERLAY_ROLES):
                continue

            # Determine severity from POLICIES
            pol_a = POLICIES[fam_a]
            pol_b = POLICIES[fam_b]
            if fam_a == fam_b:
                v = pol_a.self_overlap
            elif fam_a == Family.TEXT or fam_b == Family.TEXT:
                v = OverlapVerdict.WARN  # text over entity — role ambiguity, not physics
            else:
                v = OverlapVerdict.WARN

            lvl = _verdict_to_level(v, pol_a.strength)
            cross = ""
            if elem_region.get(a.elem_id) and elem_region.get(b.elem_id) and \
               elem_region[a.elem_id] != elem_region[b.elem_id]:
                cross = f" [cross-region: {elem_region[a.elem_id]} × {elem_region[b.elem_id]}]"
            diags.append({
                "slide": slide_idx, "phase": "2", "kind": "overlap",
                "severity": lvl, "elem_id": a.elem_id,
                "message": f"'{a.elem_id}' ({fam_a.value}) overlaps '{b.elem_id}' ({fam_b.value})"
                           f"{cross}. {v.value}.",
            })

    return diags


def _aggregate_diagnostics(diags: list[dict]) -> tuple[list[dict], dict]:
    """P0-①: collapse noisy diagnostics into a clean actionable feed.

    Rules (priority order):
      1. ALL errors pass through untouched — never deduped, batched, or trimmed.
      2. Dedup warnings: same (elem_id, kind) → keep latest phase only.
      3. Batch-collapse: >=5 warnings of same kind → one summary with elem_ids list.
      4. Cap warnings at 15, info at 5. Surplus tracked in `trimmed_*` keys.
      5. T5: signals (severity "advisory" or channel "signal") are NEVER trimmed and
         NEVER batch-folded — a focal_point.missing / image_style_conflict must always
         reach the agent even in a 20-slide deck.

    Returns (aggregated, stats) — caller MUST use aggregated as the final diagnostics list.
    """
    empty_stats = {"raw_count": 0, "errors": 0, "warnings": 0, "info": 0,
                   "advisories": 0, "dedup": 0, "batch": 0,
                   "trimmed_warnings": 0, "trimmed_info": 0, "final_count": 0}
    if not diags:
        return [], empty_stats

    errors = [d for d in diags if d.get("severity") in ("error",)]
    warns  = [d for d in diags if d.get("severity") in ("warning", "warn")]
    infos  = [d for d in diags if d.get("severity") == "info"]
    advisories = [d for d in diags
                  if d.get("severity") == "advisory" or d.get("channel") == "signal"]

    raw_counts = {"raw_count": len(diags), "errors": len(errors),
                  "warnings": len(warns), "info": len(infos),
                  "advisories": len(advisories)}

    # ── Rule 2: dedup warnings by (elem_id, kind), keep latest phase ──
    seen: dict[tuple, dict] = {}
    dedup_count = 0
    for w in warns:
        eid = w.get("elem_id", "") or ""
        kind = w.get("kind", "")
        key = (eid, kind)
        if key in seen:
            dedup_count += 1
            # Keep the one from the later phase
            existing_phase = seen[key].get("phase", "0")
            current_phase = w.get("phase", "0")
            if _phase_rank(current_phase) >= _phase_rank(existing_phase):
                seen[key] = w
        else:
            seen[key] = w
    warns = list(seen.values())

    # ── Rule 3: batch-collapse same-kind warnings (≥ threshold) ──
    by_kind: dict[str, list[dict]] = {}
    for w in warns:
        by_kind.setdefault(w.get("kind", "unknown"), []).append(w)

    batched = 0
    result_warns: list[dict] = []
    for kind, items in by_kind.items():
        if len(items) >= _KIND_BATCH_THRESHOLD:
            batched += 1
            elem_ids = sorted(set(
                i.get("elem_id", "") for i in items if i.get("elem_id")
            ))
            sample = items[0]
            result_warns.append({
                "kind": f"{kind}_batch",
                "severity": "warning",
                "phase": "aggregate",
                "batch_kind": kind,
                "count": len(items),
                "elem_ids": elem_ids,
                "sample_message": sample.get("message", ""),
                "options": sample.get("options", []),
                "message": (
                    f"{len(items)} elements share '{kind}' — "
                    f"affected: {', '.join(elem_ids[:6])}"
                    f"{'...' if len(elem_ids) > 6 else ''}"
                ),
            })
        else:
            result_warns.extend(items)

    # ── Rule 4: cap warnings and info ──
    trimmed_w = max(0, len(result_warns) - _WARN_CAP)
    result_warns = result_warns[:_WARN_CAP]

    trimmed_i = max(0, len(infos) - _INFO_CAP)
    result_infos = infos[:_INFO_CAP]

    # T5: signals are exempt from every cap/fold — appended in full, after errors+warnings.
    aggregated = errors + result_warns + result_infos + advisories
    stats = {
        **raw_counts,
        "dedup": dedup_count,
        "batch": batched,
        "trimmed_warnings": trimmed_w,
        "trimmed_info": trimmed_i,
        "final_count": len(aggregated),
    }
    return aggregated, stats


def _phase_rank(phase: str) -> int:
    """Order phases so 'keep latest phase' dedup picks the most downstream result."""
    try:
        return _PHASE_ORDER.index(phase)
    except ValueError:
        return 99  # unknown phases rank last (most recent)


def _diag(slide_idx, phase, d, kind="", severity="", deco_id="", elem_id="", message=""):
    """Normalize diagnostics from various sources into a uniform format."""
    if isinstance(d, dict):
        return {"slide": slide_idx, "phase": phase, **d}
    if hasattr(d, 'kind'):
        return {
            "slide": slide_idx, "phase": phase,
            "kind": d.kind, "severity": d.severity,
            "region_id": getattr(d, 'region_id', ''),
            "elem_id": getattr(d, 'elem_id', ''),
            "demand_pt": getattr(d, 'demand_pt', 0),
            "usable_pt": getattr(d, 'usable_pt', 0),
            "message": getattr(d, 'message', ''),
            "options": getattr(d, 'options', []),
        }
    if kind:
        return {"slide": slide_idx, "phase": phase, "kind": kind, "severity": severity,
                "deco_id": deco_id, "elem_id": elem_id, "message": message}
    return {"slide": slide_idx, "phase": phase, "message": str(d)}


def _ae_violation_to_diag(v) -> dict:
    """Convert AestheticViolation to a flat dict.

    Severity mapping: Verdict.BLOCK → "error", WARN → "warning", ALLOW → "info".
    This is canonical — BLOCK means the build must fail or the agent must explain.
    """
    severity_map = {Verdict.BLOCK: "error", Verdict.WARN: "warning", Verdict.ALLOW: "info"}
    sv = severity_map.get(v.verdict, v.verdict.name.lower() if hasattr(v.verdict, 'name') else str(v.verdict))
    return {
        "kind": v.rule_id,
        "category": v.category,
        "severity": sv,
        "priority": v.priority,
        "elem_id": v.element_id,
        "message": v.message,
        "metrics": v.metrics,
    }




# ── Per-element render hook (external streaming previews; None = off) ──
_render_frame_hook = None


def set_render_frame_hook(fn):
    """Install a per-element render callback fn(elem_id, content_type, payload, x, y, w, h),
    invoked from _render_slide right before each element is drawn (elem_id is the
    string key of canvas._phase1_rects). Returns the previous hook."""
    global _render_frame_hook
    prev = _render_frame_hook
    _render_frame_hook = fn
    return prev


def _render_slide_deco(slide, page_w_pt, page_h_pt, spec, template) -> None:
    """Slide-level decoration skins. All geometry solved here — the caller only
    declared frame/rail/corner_mark names, never coordinates."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    accent_hex = (template.accent_hex or "0052D9").lstrip("#")
    accent = RGBColor(int(accent_hex[0:2], 16), int(accent_hex[2:4], 16), int(accent_hex[4:6], 16))
    thin = 6.0  # band/rail thickness (pt)

    def bar(x, y, w, h):
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(x), Pt(y), Pt(w), Pt(h))
        sp.fill.solid()
        sp.fill.fore_color.rgb = accent
        sp.line.fill.background()
        sp.shadow.inherit = False

    if getattr(spec, "frame", "") == "top_bottom_band":
        bar(0, 0, page_w_pt, thin)
        bar(0, page_h_pt - thin, page_w_pt, thin)

    rail = getattr(spec, "rail", "")
    if rail in ("left", "right"):
        bar(0 if rail == "left" else page_w_pt - thin, 0, thin, page_h_pt)

    cm = getattr(spec, "corner_mark", "")
    if cm in ("tl", "tr"):
        bar(0 if cm == "tl" else page_w_pt - 28, 0, 28, 28)


def _render_slide(prs, canvas, template, slide_index=0, total_slides=1,
                  caption_state: dict | None = None, slide_spec=None):
    """Grid-to-PPT full slide render: background + decoration skins + info layer + page number."""
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    layout_idx = 0
    try:
        layout_count = len(prs.slide_layouts)
        if layout_count > 0:
            layout_idx = min(slide_index, layout_count - 1)
    except Exception as e:
        print(f"[PPTBuilder] slide_layouts error (template PPTX): {e}")
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    # Delete all placeholder shapes from the base layout — engine adds its own shapes.
    # Default slide layouts carry title/content placeholders that the grid engine never
    # populates, leaving ghost empty boxes in the output.
    shapes_to_delete = [s for s in slide.shapes if s.is_placeholder]
    for s in shapes_to_delete:
        sp = s._element
        sp.getparent().remove(sp)

    # Background fill
    bg_hex = template.bg_hex.lstrip("#")
    bg_rgb = RGBColor(int(bg_hex[0:2], 16), int(bg_hex[2:4], 16), int(bg_hex[4:6], 16))
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_rgb

    page_w_pt = canvas.config.canvas_w_pt
    page_h_pt = canvas.config.canvas_h_pt

    # Slide decoration skins — frame / rail / corner_mark (engine-solved geometry)
    if slide_spec is not None:
        _render_slide_deco(slide, page_w_pt, page_h_pt, slide_spec, template)

    # Phase 1: information layer elements
    for pe in canvas._phase1_rects if hasattr(canvas, '_phase1_rects') else {}:
        x, y, w, h = canvas._phase1_rects.get(pe, (0, 0, 0, 0))
        ct, payload = canvas._phase1_payloads.get(pe, (ContentType.UNKNOWN, None))
        if w <= 0 or h <= 0:
            continue
        # Per-element render hook — external streaming previews (None = off)
        if _render_frame_hook is not None:
            try:
                _render_frame_hook(pe, ct, payload, x, y, w, h)
            except Exception:
                pass

        if ct == ContentType.TABLE and payload:
            _render_table(slide, x, y, w, h, payload, template)
        elif ct == ContentType.IMAGE and payload and payload.image_path:
            # Contain-fit image rendering: PIL natural size -> contain -> no crop, no stretch
            _render_image(slide, x, y, w, h, payload)
            # Optional caption — 按 style preset 的 caption 约定渲染（前缀/对齐/字号），
            # 前缀含 "N" 时自动编号（"Figure N. " → "Figure 1. "），兑现学术主题的
            # captions_must_be_numbered 契约（旧版约定存在但从未被使用）
            if payload.caption:
                fmt = (caption_state or {}).get("format", {}) or {}
                prefix = fmt.get("prefix", "") or ""
                if prefix:
                    caption_state["n"] = caption_state.get("n", 0) + 1
                    n = caption_state["n"]
                    cap_text = (prefix.replace("N", str(n)) if "N" in prefix else prefix) + payload.caption
                else:
                    cap_text = payload.caption
                _align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
                cap_align = _align_map.get(str(fmt.get("alignment", "left")).lower(), PP_ALIGN.LEFT)
                cap_fs = float(fmt.get("font_size", 0) or template.caption_size)
                cap_x = Emu(int(x * 12700))
                cap_y = Emu(int((y + h - 6) * 12700))
                cap_w = Emu(int(w * 12700))
                cap_h = Emu(int(18 * 12700))
                tb = slide.shapes.add_textbox(cap_x, cap_y, cap_w, cap_h)
                tb.text_frame.word_wrap = True
                p = tb.text_frame.paragraphs[0]
                p.alignment = cap_align
                run = p.add_run()
                run.text = cap_text
                run.font.size = Pt(cap_fs)
                run.font.color.rgb = RGBColor(*_hex_to_rgb(template.gray_hex) if template.gray_hex else (0x66, 0x66, 0x66))
        else:
            _render_payload(slide, x, y, w, h, ct, payload,
                {"LEFT": PP_ALIGN.LEFT, "CENTER": PP_ALIGN.CENTER, "RIGHT": PP_ALIGN.RIGHT})

    # Phase 2: decorations
    for dec in (canvas._decoration_payloads if hasattr(canvas, '_decoration_payloads') else []):
        if dec.get("type") == "arrow":
            from pptx.enum.shapes import MSO_CONNECTOR_TYPE
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR_TYPE.STRAIGHT,
                Pt(dec["x1"]), Pt(dec["y1"]),
                Pt(dec["x2"]), Pt(dec["y2"]))
            connector.line.color.rgb = RGBColor(*dec.get("line_color", (0x66, 0x66, 0x66)))
            connector.line.width = Pt(dec.get("line_width_pt", 1.5))
            if dec.get("text"):
                tx = (dec["x1"] + dec["x2"]) / 2
                ty = (dec["y1"] + dec["y2"]) / 2
                from pptx.enum.text import MSO_AUTO_SIZE
                label = slide.shapes.add_textbox(Pt(tx - 55), Pt(ty - 16), Pt(110), Pt(32))
                label.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                label.text_frame.word_wrap = True
                label.text_frame.paragraphs[0].text = dec["text"]
                label.text_frame.paragraphs[0].font.size = Pt(dec.get("font_size", 10))
                label.text_frame.paragraphs[0].font.color.rgb = RGBColor(*dec.get("font_color", (0x55, 0x55, 0x55)))

    # Page number
    fn = f"{slide_index + 1}/{total_slides}"
    pn = slide.shapes.add_textbox(Pt(page_w_pt - 60), Pt(page_h_pt - 28), Pt(48), Pt(20))
    pn.text_frame.paragraphs[0].text = fn
    pn.text_frame.paragraphs[0].font.size = Pt(template.page_number_size)
    pn.text_frame.paragraphs[0].font.color.rgb = RGBColor(*_hex_to_rgb(template.dim_hex)) if template.dim_hex else RGBColor(0x88, 0x88, 0x99)
    pn.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


# ── P1-② Table rendering ──

def _render_table(slide, x: float, y: float, w: float, h: float,
                  payload, template) -> None:
    """Render a table shape into the slide. Columns auto-sized; header row styled."""
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn

    headers = payload.table_headers or []
    rows = payload.table_rows or []
    if not headers and not rows:
        return

    n_cols = len(headers) if headers else (len(rows[0]) if rows else 1)
    n_rows = 1 + len(rows)  # header + data

    # Clamp — don't crash on bad input
    n_cols = max(1, n_cols)
    n_rows = max(1, n_rows)

    col_w = w / n_cols
    row_h = min(h / n_rows, 36.0)  # cap individual row height at 36pt

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Emu(int(x * 12700)), Emu(int(y * 12700)),
        Emu(int(w * 12700)), Emu(int(h * 12700)),
    )
    tbl = tbl_shape.table

    # Column widths — equal distribution
    col_emu = Emu(int(col_w * 12700))
    for ci in range(n_cols):
        tbl.columns[ci].width = col_emu

    # Font config
    font_sz = payload.font_size if payload and payload.font_size else 12.0
    font_color = payload.font_color if payload else (0x33, 0x33, 0x44)
    header_bg = _hex_to_rgb(template.accent_hex) if template.accent_hex else (0x1A, 0x1A, 0x2E)
    if _is_dark(header_bg):
        header_font_color = (0xFF, 0xFF, 0xFF)
    else:
        header_font_color = (0xFF, 0xFF, 0xFF) if _is_dark(header_bg) else (0x33, 0x33, 0x44)

    for ri in range(n_rows):
        data = headers if ri == 0 else rows[ri - 1]
        is_header = (ri == 0)

        for ci in range(n_cols):
            cell = tbl.cell(ri, ci)
            cell_text = str(data[ci]) if ci < len(data) else ""

            # Cell fill
            cell_fill = cell.fill
            cell_fill.solid()
            if is_header:
                cell_fill.fore_color.rgb = RGBColor(*header_bg)
            else:
                bg_rgb = _hex_to_rgb(template.bg_hex) if template.bg_hex else (0xFF, 0xFF, 0xFF)
                cell_fill.fore_color.rgb = RGBColor(*bg_rgb)

            # Cell text
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(4)
            tf.margin_right = Pt(4)
            tf.margin_top = Pt(2)
            tf.margin_bottom = Pt(2)

            p = tf.paragraphs[0]
            p.text = cell_text
            p.font.size = Pt(font_sz)
            p.font.bold = is_header
            p.font.color.rgb = RGBColor(*header_font_color) if is_header else RGBColor(*font_color)

            # Vertical center
            tcPr = cell._tc.get_or_add_tcPr()
            tcPr.set('anchor', 'ctr')

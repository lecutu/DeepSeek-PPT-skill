"""
layout_extractor.py — Extract reusable slide layouts from a reference PPTX.

Bridges the gap between a human-designed .pptx and the archetype-driven
generation pipeline: reads every shape on each slide, classifies it,
normalizes its bbox to the 960×540 canvas, groups shapes into named
zones, and emits SlideLayout objects that can be registered as archetypes.

The point is NOT pixel-perfect cloning. It is structure capture: "this
reference deck has a header band, a left text column, and a right image
column" -> a layout you can fill with new content.

Usage:
    from ppt_reflex.layout_extractor import extract_layouts
    layouts = extract_layouts("company_template.pptx")

    # register as archetypes so add_slide(archetype=...) works
    from ppt_reflex.builder import PPTBuilder
    b = PPTBuilder(...)
    b.register_layouts(layouts)          # ids: "ref_0", "ref_1", ...
    b.add_slide(archetype="ref_2", ...)

Or feed one layout's regions directly:
    b.add_slide(regions=layouts[2].regions_pt, elements=[...])
"""

from __future__ import annotations
from dataclasses import dataclass, field
from typing import Optional

from ppt_reflex.grid.serializer import classify_shape, _shape_id
from ppt_reflex.grid.types import ContentType

# Canvas the engine generates on. Reference decks get normalized to this.
DEFAULT_CANVAS_W = 960.0
DEFAULT_CANVAS_H = 540.0

# ContentType → coarse semantic role used for zone naming
_BAND_TYPES = (ContentType.TEXTBOX, ContentType.SHAPE, ContentType.TABLE,
               ContentType.IMAGE, ContentType.CHART)
_TEXT_TYPES = (ContentType.TEXT, ContentType.TITLE)


@dataclass
class LayoutZone:
    """One named zone: a normalized region plus the shape semantics inside it."""
    name: str
    x: float
    y: float
    w: float
    h: float
    kind: str                 # "header" | "footer" | "band" | "text" | "image" | "deco"
    content_types: list[str] = field(default_factory=list)
    locked: bool = False      # template chrome (footer/deco) — content must not touch


@dataclass
class SlideLayout:
    """A normalized, named layout extracted from one reference slide."""
    source_slide: int
    name: str
    zones: list[LayoutZone]
    canvas_w: float = DEFAULT_CANVAS_W
    canvas_h: float = DEFAULT_CANVAS_H

    @property
    def regions_pt(self) -> list[tuple]:
        """[(name, x, y, w, h, z, inset)] — drop-in for add_slide(regions=...)."""
        out = []
        for i, z in enumerate(self.zones):
            inset = 8.0 if z.kind in ("text", "band") else 4.0
            out.append((z.name, round(z.x, 1), round(z.y, 1),
                        round(z.w, 1), round(z.h, 1), i + 1, inset))
        return out

    def to_archetype(self, archetype_id: str) -> "SlideArchetype":
        """Convert to a SlideArchetype so add_slide(archetype=...) works."""
        from ppt_reflex.grid.archetypes import SlideArchetype
        zone_map: dict[str, str] = {}
        distribute: dict[str, list[str]] = {}
        for z in self.zones:
            if z.kind == "header":
                zone_map.setdefault("title", z.name)
                zone_map.setdefault("subtitle", z.name)
            elif z.kind == "image":
                zone_map.setdefault("image", z.name)
                zone_map.setdefault("shape", z.name)
            elif z.kind == "footer":
                zone_map.setdefault("footer", z.name)
            elif z.kind in ("band", "text"):
                zone_map.setdefault("text", z.name)
                zone_map.setdefault("box", z.name)
                distribute.setdefault("box", []).append(z.name)
        # content zones that aren't single-slot: distribute text too
        content_zones = [z.name for z in self.zones
                         if z.kind in ("band", "text") and not z.locked]
        if len(content_zones) > 1:
            distribute.setdefault("text", content_zones)
        # Archetype regions are 6-tuples (name,x,y,w,h,z); add_slide injects a
        # uniform inset from LayoutPolicy. Per-zone insets live in regions_pt for
        # direct add_slide(regions=...) use.
        regions6 = [(z.name, round(z.x, 1), round(z.y, 1),
                     round(z.w, 1), round(z.h, 1), i + 1)
                    for i, z in enumerate(self.zones)]
        return SlideArchetype(
            id=archetype_id,
            name=self.name,
            description=f"Extracted from reference slide {self.source_slide}",
            regions=regions6,
            zone_map=zone_map,
            distribute=distribute,
            ai_guide=", ".join(f"{z.name}→{z.kind}" for z in self.zones),
        )


# ═══════════════════════════════════════════════════════════
# Classification helpers
# ═══════════════════════════════════════════════════════════

def _zone_kind(ct: ContentType, shape) -> str:
    if ct == ContentType.IMAGE:
        return "image"
    if ct in _TEXT_TYPES:
        return "text"
    if ct in _BAND_TYPES:
        return "band"
    return "deco"


def _is_background_fill(shape) -> bool:
    """Full-canvas backdrop — treat as chrome, not content.

    Only geometry decides: a shape covering ~90%+ of the canvas is a backdrop.
    fill.type must NOT be consulted — MSO_FILL_TYPE.BACKGROUND is python-pptx's
    "no explicit fill" (the default for text boxes), not a painted background.
    """
    try:
        w = shape.width / 12700
        h = shape.height / 12700
        if w > 860 and h > 480:
            return True
    except Exception:
        pass
    return False


def _iter_leaf_shapes(shape, dx: float = 0.0, dy: float = 0.0):
    """Walk a shape tree, yielding (leaf, abs_left_pt, abs_top_pt).

    Group children are stored relative to the group's chOff/chExt coordinate
    space; some producers (e.g. WPS) use near-absolute coordinates there, so
    the offset to add is the delta between the group's off (parent space) and
    its chOff (child space). Both are read from the raw XML because
    python-pptx reports chOff as the child's own left/top for nested groups.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            gx, gy = _group_child_offset(shape)
            for child in shape.shapes:
                yield from _iter_leaf_shapes(child, dx + gx, dy + gy)
            return
    except Exception:
        pass
    try:
        x = ((shape.left / 12700) if shape.left is not None else 0.0) + dx
        y = ((shape.top / 12700) if shape.top is not None else 0.0) + dy
    except Exception:
        return
    yield shape, x, y


def _group_child_offset(shape) -> tuple[float, float]:
    """Point offset from a group's coordinate space to its parent's.

    Standard: child coords are relative to (off - chOff). For a group whose
    off == chOff this is 0 (the common case); WPS sometimes stores near-absolute
    child coords via a large chOff, in which case off - chOff is negative and
    re-centers the child content onto the canvas.
    """
    from lxml import etree
    try:
        xml = shape._element
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        off = xml.find(".//a:off", ns)
        chOff = xml.find(".//a:chOff", ns)
        ext = xml.find(".//a:ext", ns)
        if off is not None and chOff is not None and ext is not None:
            off_x = int(off.get("x")) / 12700
            off_y = int(off.get("y")) / 12700
            ch_x = int(chOff.get("x")) / 12700
            ch_y = int(chOff.get("y")) / 12700
            return (off_x - ch_x), (off_y - ch_y)
    except Exception:
        pass
    return 0.0, 0.0


def _dedupe_zones(zones: list[LayoutZone]) -> list[LayoutZone]:
    """Merge overlapping zones that are clearly the same slot (e.g. a card
    shape + the text inside it). Keeps the container, absorbs the text.
    """
    kept: list[LayoutZone] = []
    for z in sorted(zones, key=lambda z: (z.w * z.h), reverse=True):
        absorbed = False
        for k in kept:
            if _rect_contains(k, z) or _rect_contains(z, k):
                # keep the larger container, merge semantics
                if (k.w * k.h) >= (z.w * z.h):
                    k.content_types = sorted(set(k.content_types + z.content_types))
                    if not k.locked and z.locked:
                        k.locked = z.locked
                    absorbed = True
                    break
        if not absorbed:
            kept.append(z)
    return kept


def _rect_contains(a: LayoutZone, b: LayoutZone) -> bool:
    tol = 4.0
    return (a.x - tol <= b.x and a.y - tol <= b.y and
            a.x + a.w + tol >= b.x + b.w and a.y + a.h + tol >= b.y + b.h)


def _snap_to_region(x: float, y: float, w: float, h: float) -> tuple[float, float, float, float]:
    """Snap tiny coordinate noise to a coarse grid so extracted layouts look
    intentional, not like measured 1.3pt offsets."""
    g = 16.0  # coarse grid in pt
    x0 = round(x / g) * g
    y0 = round(y / g) * g
    w0 = max(20.0, round(w / g) * g)
    h0 = max(20.0, round(h / g) * g)
    return x0, y0, w0, h0


def _name_zone(kind: str, zones: list[LayoutZone], top_y: float) -> str:
    if kind == "header":
        return "header"
    if kind == "footer":
        return "footer"
    if kind == "image":
        base = "image"
    elif kind == "text":
        base = "text"
    else:
        base = "zone"
    n = sum(1 for z in zones if z.name.startswith(base))
    return f"{base}{n + 1 if n else ''}"


# ═══════════════════════════════════════════════════════════
# Main entry
# ═══════════════════════════════════════════════════════════

def register_layouts(pptx_path: str, *, prefix: str = "ref_") -> list[str]:
    """Extract layouts from a reference deck and register each as a global
    archetype (ids: "{prefix}0", "{prefix}1", ...). Returns the registered ids.
    After this, add_slide(archetype=id) routes to the extracted layout.

    Registration is global (grid.archetypes.ARCHETYPES), so no builder changes
    are needed — get_archetype()/list_archetypes() pick them up automatically.
    """
    from ppt_reflex.grid.archetypes import register_archetype
    layouts = extract_layouts(pptx_path)
    ids: list[str] = []
    for i, layout in enumerate(layouts):
        aid = f"{prefix}{i}"
        register_archetype(layout.to_archetype(aid))
        ids.append(aid)
    return ids


def extract_layouts(pptx_path: str, *, canvas_w: float = DEFAULT_CANVAS_W,
                    canvas_h: float = DEFAULT_CANVAS_H) -> list[SlideLayout]:
    """Extract normalized layouts from every slide of a reference deck.

    Args:
        pptx_path: path to the reference .pptx
        canvas_w/canvas_h: target canvas the engine generates on. Reference
            shapes are scaled to this canvas before bbox extraction.

    Returns:
        list[SlideLayout], one per slide that has any content shapes.
    """
    from pptx import Presentation
    prs = Presentation(pptx_path)
    ref_w = prs.slide_width / 12700
    ref_h = prs.slide_height / 12700
    sx = canvas_w / ref_w if ref_w else 1.0
    sy = canvas_h / ref_h if ref_h else 1.0

    layouts: list[SlideLayout] = []
    for si, slide in enumerate(prs.slides):
        zones: list[LayoutZone] = []
        for shape in slide.shapes:
            for leaf, x0, y0 in _iter_leaf_shapes(shape):
                w = (leaf.width / 12700) * sx
                h = (leaf.height / 12700) * sy
                if w < 4 or h < 4:
                    continue
                if _is_background_fill(leaf):
                    continue
                x = x0 * sx
                y = y0 * sy

                ct = classify_shape(leaf)
                if ct == ContentType.UNKNOWN:
                    continue
                kind = _zone_kind(ct, leaf)

                # header/footer band detection by position
                if kind == "text" and y < canvas_h * 0.12 and h < 60:
                    kind = "header"
                elif kind in ("text", "band") and y > canvas_h * 0.88 and h < 60:
                    kind = "footer"

                locked = kind in ("header", "footer", "deco")
                zones.append(LayoutZone(
                    name="",
                    x=round(x, 1), y=round(y, 1), w=round(w, 1), h=round(h, 1),
                    kind=kind,
                    content_types=[ct.value],
                    locked=locked,
                ))

        zones = _dedupe_zones(zones)
        zones = [z for z in zones if not (z.kind == "deco") or z.locked]
        if not zones:
            continue

        for z in zones:
            z.x, z.y, z.w, z.h = _snap_to_region(z.x, z.y, z.w, z.h)
            z.name = _name_zone(z.kind, zones, z.y)

        layouts.append(SlideLayout(
            source_slide=si,
            name=f"Layout {si + 1}",
            zones=zones,
            canvas_w=canvas_w,
            canvas_h=canvas_h,
        ))
    return layouts


def summary(layouts: list[SlideLayout]) -> dict:
    """Human-readable summary for the caller / diagnostics."""
    return {
        "layouts": len(layouts),
        "detail": [
            {
                "index": i,
                "source_slide": l.source_slide,
                "name": l.name,
                "zones": [
                    {"name": z.name, "kind": z.kind, "xywh": (z.x, z.y, z.w, z.h),
                     "locked": z.locked}
                    for z in l.zones
                ],
            }
            for i, l in enumerate(layouts)
        ],
    }
